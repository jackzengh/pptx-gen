"""Shared building blocks for the check_pptx layout validator.

This module holds everything that is NOT a rubric check: the unit conversions,
the in-memory data model (Box / SlideInfo / DeckContext), the pipeline that
flattens a python-pptx slide into absolute-coordinate Boxes, the role/kind
classifiers, and the small geometry/clustering helpers the checks lean on.

`tools.py` imports these and implements the check_* rules on top. Keeping them
apart means the checks read as pure policy ("a row's peers must share a width")
while all the fiddly EMU math and XML traversal lives here, justified in place.
"""
from __future__ import annotations

from dataclasses import dataclass, field

from pptx.enum.shapes import MSO_SHAPE_TYPE, PP_PLACEHOLDER
from pptx.oxml.ns import qn

EMU_PER_INCH = 914400

# TOL is the single alignment/overlap tolerance, in inches. Two coordinates
# within TOL count as "the same" — this absorbs the sub-0.1in rounding that
# real authoring tools introduce, so a 3-column layout whose tops differ by
# 0.02in isn't flagged as misaligned. Every check shares this one value so the
# notion of "aligned" is consistent across the whole validator.
TOL = 0.1

# MARGIN_BAND is the reserved outer margin (inches). It is deliberately MUCH
# larger than TOL: TOL is "are these two things aligned?", MARGIN_BAND is "is
# this shape too close to the slide edge?". 0.3in matches the rubric grader's
# band and standard investment-banking slide convention, keeping page numbers /
# footers / content clear of the trim edge.
MARGIN_BAND = 0.3

# Max vertical gap (inches) between a chart/table and the caption above it for
# the caption to "belong" to that element (check_element_boxing). Anything
# farther reads as an unrelated label, not a caption.
CAPTION_GAP = 0.5


def emu_to_in(value: int) -> float:
    inches = value.inches if hasattr(value, "inches") else value / EMU_PER_INCH
    return round(inches, 3)


# --------------------------------------------------------------------------- #
# Data model
# --------------------------------------------------------------------------- #

@dataclass
class Box:
    # characteristics of a shape (whether its text or a shape)

    name: str
    kind: str          # title | body | chart | table | picture | text | other
    left: float
    top: float
    width: float
    height: float
    text: str = ""
    
    # Paragraph only values
    paras: list = field(default_factory=list)
    text_insets: tuple = (0.1, 0.1, 0.05, 0.05)  # (left, right, top, bottom)

    @property
    def right(self) -> float:
        return round(self.left + self.width, 3)

    @property
    def bottom(self) -> float:
        return round(self.top + self.height, 3)

    def as_dict(self) -> dict:
        return {
            "left": self.left,
            "top": self.top,
            "width": self.width,
            "height": self.height,
        }


@dataclass
class SlideInfo:
    # info about eaach slide which contains boxes

    index: int                       # 1-based, for human-facing messages
    boxes: list[Box] = field(default_factory=list)
    notes: list[str] = field(default_factory=list)


@dataclass
class DeckContext:
    # info about the whole deck which contains slides, and gives slide size (consistent across all slides)

    slide_w: float                   # inches
    slide_h: float                   # inches
    slides: list[SlideInfo]



def role_of(b: Box):
    # layout role read from the shape's name, or None for plain content.
    # tells the checks which shapes to skip and why:
    #   background  full-bleed panel/image — meant to touch the edge, content sits on it
    #   chartpart   a piece of a hand-drawn chart; two chartparts may overlap (one unit)
    #   overlay     a label sitting on something (caption/badge) — exempt from overlap
    #   meta        non-grid chrome like legend, page numbers, footers
    n = (b.name or "").lower()
    for role in ("background", "chartpart", "overlay", "meta"):
        if n == role or n.startswith(role + "_") or n.endswith("_" + role):
            return role
    return None


def kind_of(shape, title_shape) -> str:
    # classify the type of shape - title, chart, etc.

    if shape is title_shape:
        return "title"
    
    # check if the shape has a table or chart
    if getattr(shape, "has_table", False):
        return "table"
    if getattr(shape, "has_chart", False):
        return "chart"
    if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
        return "picture"
    if shape.is_placeholder:
        ph_type = shape.placeholder_format.type
        if ph_type is PP_PLACEHOLDER.TITLE or ph_type is PP_PLACEHOLDER.CENTER_TITLE:
            return "title"
        if ph_type is PP_PLACEHOLDER.BODY:
            return "body"
    if (getattr(shape, "has_text_frame", False) and shape.has_text_frame
            and shape.text_frame.text.strip()):
        return "text"
    return "other"


def is_child(child: Box, parent: Box, pad: float = TOL) -> bool:
    # is the element something inside of another box (like rectange inside of a rectangle for aesthetic purposes)
    
    return (child.left >= parent.left - pad and child.top >= parent.top - pad
            and child.right <= parent.right + pad and child.bottom <= parent.bottom + pad)


def grid_boxes(slide: SlideInfo) -> list[Box]:
    # which shapes need to abide by alignemnt,etc. ? Some shapes don't because they're inside of other shpaes, they're backgrounds ,etc.
    panels = [b for b in slide.boxes if role_of(b) == "background"] # backgrounds are excluded from the grid
    out = []
    for b in slide.boxes:
        if role_of(b) is not None:           
            continue
        if any(is_child(b, p) for p in panels):  # its a child
            continue
        out.append(b)
    return out


# --------------------------------------------------------------------------- #
# Shape collection: flatten a python-pptx slide into absolute-coord Boxes
# --------------------------------------------------------------------------- #

def _group_transform(group):
    # maps a child's box into parent coords. children are positioned relative
    # to the group's origin, but we need their position on the slide as a whole.
    xfrm = group._element.grpSpPr.xfrm
    ox, oy = xfrm.off.x, xfrm.off.y
    ecx, ecy = xfrm.ext.cx, xfrm.ext.cy
    cox, coy = xfrm.chOff.x, xfrm.chOff.y
    chcx, chcy = xfrm.chExt.cx, xfrm.chExt.cy

    sx = (ecx / chcx) if chcx else 1.0
    sy = (ecy / chcy) if chcy else 1.0

    def to_parent(left, top, width, height):
        return (
            ox + (left - cox) * sx,
            oy + (top - coy) * sy,
            width * sx,
            height * sy,
        )

    return to_parent


def _text_metrics(shape):
    
    # paragraphs are the text in the shape
    paras = []
    # insets are the text-frame's internal padding (l/r/t/b, inches)
    insets = (0.1, 0.1, 0.05, 0.05)
    if not (getattr(shape, "has_text_frame", False) and shape.has_text_frame):
        return paras, insets

    tf = shape.text_frame 

    # whats the size, bold, italic, and name of the text in the shape?
    for p in tf.paragraphs:
        size = p.font.size.pt if p.font.size is not None else None
        bold = bool(p.font.bold)
        italic = bool(p.font.italic)
        name = p.font.name

        for run in p.runs:
            if run.font.size is not None:
                size = run.font.size.pt
            bold = bold or bool(run.font.bold)
            italic = italic or bool(run.font.italic)
            if run.font.name:
                name = run.font.name

        paras.append((p.text, size, bold, italic, name))

    # get the text-frame's internal padding (l/r/t/b, inches)
    body_pr = tf._txBody.find(qn("a:bodyPr"))
    if body_pr is not None:
        def inset(attr, default):
            value = body_pr.get(attr)
            return int(value) / EMU_PER_INCH if value is not None else default

        insets = (
            inset("lIns", 0.1), inset("rIns", 0.1),
            inset("tIns", 0.05), inset("bIns", 0.05),
        )
    # returns the paragraphs and insets in inches
    return paras, insets


def _walk(shape, title_shape, transforms, out: list[Box], notes: list[str]):
    
    # get all of the shapes in the group and their cooridnates
    if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
        to_parent = _group_transform(shape)
        for child in shape.shapes:
            _walk(child, None, transforms + [to_parent], out, notes)
        return

    # get the coordinates of the shape
    left, top = shape.left, shape.top
    width, height = shape.width, shape.height
    if None in (left, top, width, height): # if the shape has no coordinates, skip it
        notes.append(f"{shape.name or 'shape'}: inherited geometry (skipped)")
        return

    # convert the coordinates to inches
    left, top, width, height = int(left), int(top), int(width), int(height)
    for to_parent in reversed(transforms):
        left, top, width, height = to_parent(left, top, width, height)

    has_text = getattr(shape, "has_text_frame", False) and shape.has_text_frame # check if the shape has text
    paras, insets = _text_metrics(shape)

    # create a box for the shape
    out.append(
        Box(
            name=shape.name or "shape",
            kind=kind_of(shape, title_shape),
            left=emu_to_in(left),
            top=emu_to_in(top),
            width=emu_to_in(width),
            height=emu_to_in(height),
            text=shape.text_frame.text.strip() if has_text else "",
            paras=paras,
            text_insets=insets,
        )
    )


def collect_shapes(slide, index: int) -> SlideInfo:
    # flatten a slide into absolute-coordinate Boxes; _walk un-nests grouped shapes.
    info = SlideInfo(index=index + 1)
    title_shape = slide.shapes.title  # may be None
    for shape in slide.shapes:
        _walk(shape, title_shape, [], info.boxes, info.notes)
    return info

def overlap_ink_intersection_area(a: Box, b: Box) -> float:
    # overlap area (in²) of two boxes' declared rectangles; 0 if disjoint.
    dx = min(a.right, b.right) - max(a.left, b.left)
    dy = min(a.bottom, b.bottom) - max(a.top, b.top)
    return dx * dy if dx > 0 and dy > 0 else 0.0


# --------------------------------------------------------------------------- #
# check how thingos are grouped togehter! 
# --------------------------------------------------------------------------- #

def alignment_cluster_1d(values: list[float]) -> list[list[float]]:
    # greedy 1-D clustering — folds many near-equal edges into a few distinct
    # "margins"; a value joins the first cluster within TOL of its running mean.
    clusters: list[list[float]] = []
    for v in sorted(values):
        for c in clusters:
            if abs(v - (sum(c) / len(c))) <= TOL:
                c.append(v)
                break
        else:
            clusters.append([v])
    return clusters


def alignment_peer_clusters(boxes: list[Box], primary: str, secondary: str) -> list[list[Box]]:
    # group peer boxes sharing a primary edge + secondary size — one helper for
    # both rows ("top","height") and columns ("left","width").
    peers = [b for b in boxes if b.kind in {"body", "chart", "table", "picture", "text"}]
    groups: list[list[Box]] = []
    for b in sorted(peers, key=lambda x: getattr(x, primary)):
        for g in groups:
            if (abs(getattr(b, primary) - getattr(g[0], primary)) <= TOL
                    and abs(getattr(b, secondary) - getattr(g[0], secondary)) <= TOL):
                g.append(b)
                break
        else:
            groups.append([b])
    return [g for g in groups if len(g) >= 2]


def alignment_block_left_edges(boxes: list[Box]) -> list[float]:
    # one left edge per block-row, so a clean N-column layout reads as one margin
    # not N; collapses strict peer rows and "soft rows" sharing only a top.
    peer_rows = alignment_peer_clusters(boxes, "top", "height")
    row_members = {id(b) for row in peer_rows for b in row}
    lefts = [min(b.left for b in row) for row in peer_rows]

    standalone = [
        b for b in boxes
        if b.kind in {"title", "body", "chart", "table", "picture", "text"}
        and id(b) not in row_members
    ]
    soft_rows: list[list[Box]] = []
    for b in sorted(standalone, key=lambda x: x.top):
        for r in soft_rows:
            if abs(b.top - r[0].top) <= TOL:
                r.append(b)
                break
        else:
            soft_rows.append([b])
    lefts.extend(min(b.left for b in r) for r in soft_rows)
    return lefts


def y_offset_variation(label, tops: dict[int, float]) -> list[dict]:
    # flag when a fixed element (title/footer) drifts in y across the deck —
    # given {slide: top_y}, emit one problem if the spread exceeds TOL.
    if len(tops) < 2:
        return []
    values = list(tops.values())
    if max(values) - min(values) > TOL:
        slides = sorted(tops.keys())
        return [{
            "slide": slides[0],
            "message": f"{label} top y-offset varies across slides "
            f"({round(min(values), 3)}..{round(max(values), 3)} in)",
            "shapes": [],
            "boxes_in": [],
        }]
    return []
