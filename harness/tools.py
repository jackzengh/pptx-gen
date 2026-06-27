from __future__ import annotations

from itertools import combinations

from pptx import Presentation

from .text_metrics import measure_text_emu
from .utils import (
    DeckContext,
    TOL, MARGIN_BAND, CAPTION_GAP, EMU_PER_INCH,
    role_of, is_child, grid_boxes,
    collect_shapes, emu_to_in,
    overlap_ink_intersection_area,
    alignment_cluster_1d, alignment_peer_clusters, alignment_block_left_edges,
    y_offset_variation,
)

# Per-line vertical pitch (single-spaced ≈ 1.2x point size) for the text-overflow height estimate.
LINE_SPACING = 1.2

# Default point sizes per kind, in case we can't get them from the shape (like if its set by the parent or the slide)
DEFAULT_PT = {"title": 36.0, "body": 14.0, "text": 14.0}


def check_outer_margin_frame(ctx: DeckContext) -> list[dict]:
    # every shape must stay inside a 0.3in frame on all 4 edges, except
    # background panels/images which are meant to bleed to the edge. We use role_of to check if the shape is a background! 
    out = []
    for s in ctx.slides:
        for b in s.boxes:
            # Only full-bleed panels / image fields are meant to touch the edge.
            # (A caption that must reach the edge sits over a background_ field.)
            if role_of(b) == "background":
                continue
            off_slide = (
                b.left < 0 or b.top < 0
                or b.right > ctx.slide_w or b.bottom > ctx.slide_h
            )
            in_band = (
                b.left < MARGIN_BAND or b.top < MARGIN_BAND
                or b.right > ctx.slide_w - MARGIN_BAND
                or b.bottom > ctx.slide_h - MARGIN_BAND
            )
            if off_slide:
                out.append({
                    "slide": s.index,
                    "message": f"{b.name!r} bleeds off the slide",
                    "shapes": [b.name],
                    "boxes_in": [b.as_dict()],
                })
            elif in_band:
                out.append({
                    "slide": s.index,
                    "message": f"{b.name!r} intrudes into the {MARGIN_BAND} in outer margin band",
                    "shapes": [b.name],
                    "boxes_in": [b.as_dict()],
                })
    return out


def check_no_overlap(ctx: DeckContext) -> list[dict]:
    # no two shapes' should collide, but some overlaps are intentional (like backgrounds!). We exclude those using 'roles' which we use the role_of function to check.
    out = []
    for s in ctx.slides:
        for a, b in combinations(s.boxes, 2):
            area = overlap_ink_intersection_area(a, b)
            if area <= 0:
                continue
            role_a = role_of(a)
            role_b = role_of(b)
            if "background" in (role_a, role_b) or "overlay" in (role_a, role_b):
                continue
            if role_a == "chartpart" and role_b == "chartpart":
                continue
            out.append({
                "slide": s.index,
                "message": f"{a.name!r} and {b.name!r} overlap by "
                f"{round(area, 3)} in^2",
                "shapes": [a.name, b.name],
                "boxes_in": [a.as_dict(), b.as_dict()],
            })
    return out


def check_left_margin_consistency(ctx: DeckContext) -> list[dict]:
    # content should hang off one shared left edge, per slide and deck-wide. so basic left margin consistency 
    out = []

    slide_lefts = {}
    for s in ctx.slides:
        lefts = alignment_block_left_edges(grid_boxes(s))
        if not lefts:
            continue
        clusters = alignment_cluster_1d(lefts)
        if len(clusters) > 2:
            out.append({
                "slide": s.index,
                "message": f"slide uses {len(clusters)} distinct block left edges "
                f"(expected at most 2)",
                "shapes": [],
                "boxes_in": [],
            })
        dominant = max(clusters, key=lambda c: len(c))
        slide_lefts[s.index] = sum(dominant) / len(dominant)

    if len(slide_lefts) >= 2:
        # Deck dominant left = most common slide-dominant value (within TOL).
        values = list(slide_lefts.values())
        deck_clusters = alignment_cluster_1d(values)
        deck_dominant_cluster = max(deck_clusters, key=lambda c: len(c))
        deck_left = sum(deck_dominant_cluster) / len(deck_dominant_cluster)
        deviating = [
            idx for idx, v in slide_lefts.items()
            if abs(v - deck_left) > TOL
        ]
        if len(deviating) >= 2:
            out.append({
                "slide": deviating[0],
                "message": f"{len(deviating)} slides (e.g. {sorted(deviating)}) use a left "
                f"edge differing by more than {TOL} in from the deck's dominant "
                f"left edge of {round(deck_left, 3)} in",
                "shapes": [],
                "boxes_in": [],
            })
    return out


def check_edge_and_grid_alignment(ctx: DeckContext) -> list[dict]:
    # peer boxes in a row must share a width and a top (a clean grid band).
    out = []
    for s in ctx.slides:
        for row in alignment_peer_clusters(grid_boxes(s), "top", "height"):
            if len(row) < 2:
                continue
            widths = [b.width for b in row]
            tops = [b.top for b in row]
            if max(widths) - min(widths) > TOL:
                out.append({
                    "slide": s.index,
                    "message": f"peer boxes in a row have unequal widths "
                    f"({round(min(widths), 3)}..{round(max(widths), 3)} in)",
                    "shapes": [b.name for b in row],
                    "boxes_in": [b.as_dict() for b in row],
                })
            if max(tops) - min(tops) > TOL:
                out.append({
                    "slide": s.index,
                    "message": f"peer boxes in a row are not top-aligned "
                    f"({round(min(tops), 3)}..{round(max(tops), 3)} in)",
                    "shapes": [b.name for b in row],
                    "boxes_in": [b.as_dict() for b in row],
                })
    return out


def check_consistent_gutters(ctx: DeckContext) -> list[dict]:
    # gaps between 3+ peers in a row or column must be even. cehcks that they do
    out = []
    for s in ctx.slides:
        # horizontal gutters within rows of 3+
        for row in alignment_peer_clusters(grid_boxes(s), "top", "height"):
            if len(row) < 3:
                continue
            ordered = sorted(row, key=lambda b: b.left)
            gaps = [
                ordered[i + 1].left - ordered[i].right
                for i in range(len(ordered) - 1)
            ]
            if gaps and max(gaps) - min(gaps) > TOL:
                out.append({
                    "slide": s.index,
                    "message": f"horizontal gutters in a row of {len(row)} vary "
                    f"({round(min(gaps), 3)}..{round(max(gaps), 3)} in)",
                    "shapes": [b.name for b in ordered],
                    "boxes_in": [b.as_dict() for b in ordered],
                })
        # vertical gutters within columns of 3+ — drop titles and bottom-band
        # footer/source boxes first (distinct roles, not list peers)
        for col in alignment_peer_clusters(grid_boxes(s), "left", "width"):
            peers = [
                b for b in col
                if b.kind != "title" and b.bottom < ctx.slide_h - 0.6
            ]
            if len(peers) < 3:
                continue
            ordered = sorted(peers, key=lambda b: b.top)
            gaps = [
                ordered[i + 1].top - ordered[i].bottom
                for i in range(len(ordered) - 1)
            ]
            if gaps and max(gaps) - min(gaps) > TOL:
                out.append({
                    "slide": s.index,
                    "message": f"vertical gutters in a column of {len(col)} vary "
                    f"({round(min(gaps), 3)}..{round(max(gaps), 3)} in)",
                    "shapes": [b.name for b in ordered],
                    "boxes_in": [b.as_dict() for b in ordered],
                })
    return out


def check_element_boxing(ctx: DeckContext) -> list[dict]:
    # every chart/table needs a caption directly above it, within CAPTION_GAP above and horizontally overlapping it. Checke that this exists
    out = []
    for s in ctx.slides:
        captions = [b for b in s.boxes if b.kind in ("text", "title") and b.text]
        for el in s.boxes:
            if el.kind not in ("chart", "table"):
                continue
            has_caption = any(
                # caption sits above the element, within CAPTION_GAP, x-overlapping
                (el.top - cap.bottom) >= -TOL
                and (el.top - cap.bottom) <= CAPTION_GAP
                and min(cap.right, el.right) - max(cap.left, el.left) > 0
                for cap in captions
            )
            if not has_caption:
                out.append({
                    "slide": s.index,
                    "message": f"{el.kind} {el.name!r} has no caption within "
                    f"{CAPTION_GAP} in above it",
                    "shapes": [el.name],
                    "boxes_in": [el.as_dict()],
                })
    return out


def check_text_overflow(ctx: DeckContext) -> list[dict]:
    # we have no renderer, so we estimate text height using a proxy font called Liberation Sans
    out = []
    for s in ctx.slides:
        for b in s.boxes:
            # measure the rendered text height (inches) of this box, inline.
            if not b.paras or all(not (txt or "").strip() for txt, *_ in b.paras):
                continue  # no measurable text
            left_in, right_in, top_in, bot_in = b.text_insets
            usable_w_in = b.width - left_in - right_in
            if usable_w_in <= 0:
                continue
            box_w_emu = int(usable_w_in * EMU_PER_INCH)
            default_pt = DEFAULT_PT.get(b.kind, 14.0)
            est = top_in + bot_in
            for text, size, bold, italic, name in b.paras:
                _, h_emu = measure_text_emu(
                    text or "",
                    box_w_emu,
                    size if size is not None else default_pt,
                    LINE_SPACING,
                    font_family=name or "Arial",
                    bold=bool(bold),
                    italic=bool(italic),
                )
                est += h_emu / EMU_PER_INCH

            text_bottom = b.top + est
            if text_bottom <= b.bottom + TOL:
                continue  # text fits inside its own box
            # Overspill exists: does it land on a shape below?
            for other in s.boxes:
                if other is b:
                    continue
                # horizontal overlap with the text box
                if min(b.right, other.right) - max(b.left, other.left) <= 0:
                    continue
                # vertical overlap of the overspill band with the other shape
                overlap = (min(text_bottom, other.bottom)
                           - max(b.bottom, other.top))
                if overlap > TOL:
                    out.append({
                        "slide": s.index,
                        "message": f"{b.name!r} text overflows its box (≈{round(est, 2)} in "
                        f"tall vs {round(b.height, 2)} in) and collides with "
                        f"{other.name!r} below",
                        "shapes": [b.name, other.name],
                        "boxes_in": [b.as_dict(), other.as_dict()],
                    })
                    break
    return out


def check_y_offset_variation(ctx: DeckContext) -> list[dict]:
    # title and footer must sit at one fixed y across the deck.
    # which titles count? grid titles only ignore the rest :))
    out = []
    title_tops = {}
    footer_tops = {}
    for s in ctx.slides:
        grid = grid_boxes(s)
        panels = [b for b in s.boxes if role_of(b) == "background"]
        titles = [b for b in grid if b.kind == "title"]
        if titles:
            title_tops[s.index] = min(b.top for b in titles)
        # Footer/source rhythm: only the standard bottom-band source line, NOT a
        # footer placed inside a sidebar panel (cover/section slides).
        footers = [
            b for b in s.boxes
            if b.kind in ("text", "body") and b.top > ctx.slide_h - 1.0
            and not any(is_child(b, p) for p in panels)
        ]
        if footers:
            footer_tops[s.index] = max(footers, key=lambda b: b.top).top

    out += y_offset_variation("title", title_tops)
    out += y_offset_variation("footer/page-number", footer_tops)
    return out

def check_pptx(path: str) -> dict:
    # main entry: validate one .pptx and return {ok, problems, notes}. aggs all the checks tgt! 
    prs = Presentation(path)
    slides = [collect_shapes(s, i) for i, s in enumerate(prs.slides)]
    ctx = DeckContext(
        slide_w=emu_to_in(prs.slide_width),
        slide_h=emu_to_in(prs.slide_height),
        slides=slides,
    )

    problems: list[dict] = []
    problems += check_outer_margin_frame(ctx)
    problems += check_no_overlap(ctx)
    problems += check_left_margin_consistency(ctx)
    problems += check_edge_and_grid_alignment(ctx)
    problems += check_consistent_gutters(ctx)
    problems += check_element_boxing(ctx)
    problems += check_text_overflow(ctx)
    problems += check_y_offset_variation(ctx)

    notes = [
        {"slide": s.index, "notes": s.notes}
        for s in slides if s.notes
    ]
    return {
        "ok": not problems,
        "problems": problems,
        "notes": notes,
    }
