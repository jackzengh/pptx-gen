#!/usr/bin/env python3
"""Exemplar: build the Take-Two financial overview deck with python-pptx.

Self-contained reference for the `pptx` skill. Every element is created through
the inlined helpers (textbox / paragraph / rectangle / connector / title / chart
/ table) so the whole deck satisfies check_pptx and the design rules:
  - action-title sentences, each with an accent connector beneath (title())
  - auto_size=SHAPE_TO_FIT_TEXT on every text frame (textbox())
  - real p.level bullets with hanging indent (paragraph(bullet=True))
  - charts never show the default "Chart Title" (chart())
  - readable tables (table())
  - true <p:bg> full-bleed backgrounds (not shapes)
  - a fixed top-right motif on every content slide
Run from inside repo/ (title() imports harness.text_metrics):
  PYTHONPATH=<repo-parent> python build_pptx.py
"""
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR, MSO_AUTO_SIZE
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.chart.data import CategoryChartData
from pptx.oxml.ns import qn
from pptx.util import Emu, Pt, Inches

# Real-font text measurement (bundled Liberation faces) — used to size titles
# for their actual wrap so connectors/content never strike through wrapped text.
from harness.text_metrics import measure_text_emu

# --------------------------------------------------------------------------- #
# Helpers (inlined from the pptx skill so this file is copy-pasteable)
# --------------------------------------------------------------------------- #
EMU = 914400
_EMU_PER_POINT = EMU / 72
_BULLET_GLYPHS = ["•", "–", "◦"]
_HANG_EM, _STEP_EM = 1.1, 1.4

MARGIN_IN = 0.6
TITLE_TOP_IN = 0.5
TITLE_PT = 30


def bullet_indent_emu(level, size_pt):
    em = size_pt * _EMU_PER_POINT
    return round((_HANG_EM + _STEP_EM * max(0, level)) * em), round(-_HANG_EM * em)


def _apply_bullet(p, level, size_pt, font):
    marL, indent = bullet_indent_emu(level, size_pt)
    pPr = p._p.get_or_add_pPr()
    pPr.set("marL", str(marL)); pPr.set("indent", str(indent))
    for tag in ("a:buNone", "a:buChar", "a:buAutoNum", "a:buFont"):
        for el in pPr.findall(qn(tag)):
            pPr.remove(el)
    g = _BULLET_GLYPHS[min(max(0, level), len(_BULLET_GLYPHS) - 1)]
    pPr.append(pPr.makeelement(qn("a:buFont"), {"typeface": font}))
    pPr.append(pPr.makeelement(qn("a:buChar"), {"char": g}))


def style_run(run, size_pt, color, *, bold=False, italic=False, font="Arial"):
    run.font.size = Pt(size_pt); run.font.color.rgb = color
    run.font.bold = bold; run.font.italic = italic; run.font.name = font


def textbox(slide, left, top, width, height, *, name):
    tb = slide.shapes.add_textbox(Emu(int(left)), Emu(int(top)), Emu(int(width)), Emu(int(height)))
    tb.name = name
    tf = tb.text_frame
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT
    tf.margin_left = tf.margin_right = Pt(0)
    tf.margin_top = tf.margin_bottom = Pt(0)
    return tb, tf


def paragraph(tf, text, size_pt, color, *, bold=False, italic=False, level=0,
              align=PP_ALIGN.LEFT, space_after_pt=6, space_before_pt=0,
              first=False, font="Arial", bullet=False):
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    p.level = level; p.alignment = align
    p.space_after = Pt(space_after_pt); p.space_before = Pt(space_before_pt)
    if bullet:
        _apply_bullet(p, level, size_pt, font)
    run = p.add_run(); run.text = text
    style_run(run, size_pt, color, bold=bold, italic=italic, font=font)
    return p


def rectangle(slide, left, top, width, height, *, name, fill=None, line=None, line_w_pt=0.75):
    sp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Emu(int(left)), Emu(int(top)),
                                Emu(int(width)), Emu(int(height)))
    sp.name = name
    sp.shadow.inherit = False
    if fill is None:
        sp.fill.background()
    else:
        sp.fill.solid(); sp.fill.fore_color.rgb = fill
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line; sp.line.width = Pt(line_w_pt)
    return sp


def connector(slide, left, top, width, color, *, name, weight_pt=1.5):
    ln = slide.shapes.add_connector(2, Emu(int(left)), Emu(int(top)), Emu(int(left + width)), Emu(int(top)))
    ln.name = name
    ln.line.color.rgb = color; ln.line.width = Pt(weight_pt)
    return ln


def title(slide, text, *, accent, ink, width_in=11.0, font="Georgia",
          subtitle=None, muted=None, name="title"):
    from harness.text_metrics import measure_text_emu
    left = int(MARGIN_IN * EMU); top = int(TITLE_TOP_IN * EMU); w = int(width_in * EMU)
    _, h_emu = measure_text_emu(text, w, TITLE_PT, 1.15, font_family=font, bold=True)
    h_emu = max(h_emu, int(0.7 * EMU))
    _, tf = textbox(slide, left, top, w, h_emu, name=name)
    paragraph(tf, text, TITLE_PT, ink, bold=True, first=True, font=font, space_after_pt=0)
    conn_y = top + h_emu + int(0.10 * EMU)
    connector(slide, left, conn_y, int(1.6 * EMU), accent, name=f"{name}_rule", weight_pt=2.5)
    content_top = conn_y + int(0.12 * EMU)
    if subtitle:
        st_y = content_top
        _, stf = textbox(slide, left, st_y, w, int(0.3 * EMU), name=f"caption_{name}")
        paragraph(stf, subtitle, 11, muted or ink, first=True, font="Arial")
        content_top = st_y + int(0.34 * EMU)
    return content_top


def chart(slide, left, top, width, height, categories, series, *, palette,
          name, kind="column", legend_bottom=True):
    data = CategoryChartData()
    data.categories = categories
    for label, vals in series:
        data.add_series(label, vals)
    xl = {"column": XL_CHART_TYPE.COLUMN_CLUSTERED, "bar": XL_CHART_TYPE.BAR_CLUSTERED,
          "line": XL_CHART_TYPE.LINE_MARKERS, "pie": XL_CHART_TYPE.PIE}[kind]
    gf = slide.shapes.add_chart(xl, Emu(int(left)), Emu(int(top)), Emu(int(width)), Emu(int(height)), data)
    gf.name = name
    ch = gf.chart
    ch.has_title = False
    if kind == "pie":
        # A pie has one series; color each POINT (slice) from the palette and
        # show the legend so slices are labeled.
        ch.has_legend = True
        ch.legend.position = XL_LEGEND_POSITION.BOTTOM
        ch.legend.include_in_layout = False
        pts = ch.plots[0].series[0].points
        for i, pt in enumerate(pts):
            pt.format.fill.solid()
            pt.format.fill.fore_color.rgb = palette[i % len(palette)]
        return gf
    ch.has_legend = bool(legend_bottom) and len(series) > 1
    if ch.has_legend:
        ch.legend.position = XL_LEGEND_POSITION.BOTTOM
        ch.legend.include_in_layout = False
    for i, ps in enumerate(ch.plots[0].series):
        ps.format.fill.solid()
        ps.format.fill.fore_color.rgb = palette[i % len(palette)]
    return gf


def table(slide, left, top, width, height, headers, rows, *, header_fill, header_text,
          ink, name, col_widths_in=None, header_pt=12, body_pt=11, row_h_in=0.4):
    gf = slide.shapes.add_table(len(rows) + 1, len(headers), Emu(int(left)), Emu(int(top)),
                                Emu(int(width)), Emu(int(height)))
    gf.name = name
    t = gf.table
    for j, htext in enumerate(headers):
        c = t.cell(0, j)
        c.fill.solid(); c.fill.fore_color.rgb = header_fill
        c.text = str(htext)
        p = c.text_frame.paragraphs[0]
        p.runs[0].font.size = Pt(header_pt); p.runs[0].font.bold = True
        p.runs[0].font.color.rgb = header_text
    for i, row in enumerate(rows, start=1):
        for j, val in enumerate(row):
            c = t.cell(i, j)
            c.text = str(val)
            r = c.text_frame.paragraphs[0].runs[0]
            r.font.size = Pt(body_pt); r.font.color.rgb = ink
    for i in range(len(rows) + 1):
        t.rows[i].height = Emu(int(row_h_in * EMU))
    if col_widths_in:
        for j, w_in in enumerate(col_widths_in):
            t.columns[j].width = Emu(int(w_in * EMU))
    return gf


# --------------------------------------------------------------------------- #
# Palette + deck-level constants
# --------------------------------------------------------------------------- #
INK = RGBColor(0x1B, 0x26, 0x38)
MUTED = RGBColor(0x5B, 0x66, 0x77)
FAINT = RGBColor(0x8A, 0x93, 0xA1)
AMBER = RGBColor(0xE3, 0x6A, 0x1E)
BLUE = RGBColor(0x2E, 0x5E, 0xAA)
TEAL = RGBColor(0x2A, 0x9D, 0x8F)
CREAM = RGBColor(0xF7, 0xF3, 0xEC)
CARD = RGBColor(0xFD, 0xFB, 0xF6)
CARD_LINE = RGBColor(0xE0, 0xD8, 0xC9)
DARK = RGBColor(0x14, 0x1C, 0x2B)
LIGHTTX = RGBColor(0xE9, 0xED, 0xF3)

SW, SH = 13.333, 7.5
M = int(MARGIN_IN * EMU)            # left margin EMU
RIGHT = int((SW - MARGIN_IN) * EMU)  # right content edge EMU
CHIPS = [AMBER, BLUE, TEAL]
SRC_Y = int(7.05 * EMU)             # fixed source-line y (clear of 0.1in bottom band)

prs = Presentation()
prs.slide_width = Inches(SW)
prs.slide_height = Inches(SH)
BLANK = prs.slide_layouts[6]


def new_slide(bg):
    s = prs.slides.add_slide(BLANK)
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = bg
    return s


def motif(slide):
    """Fixed top-right three-chip motif (same position on every content slide)."""
    cw, ch, gap = int(0.22 * EMU), int(0.1 * EMU), int(0.1 * EMU)
    top = int(0.55 * EMU)
    x = RIGHT - 3 * cw - 2 * gap
    for i, c in enumerate(CHIPS):
        rectangle(slide, x + i * (cw + gap), top, cw, ch, name=f"motif_chip_{i}", fill=c)


def source(slide, text, *, color=FAINT):
    _, tf = textbox(slide, M, SRC_Y, int((SW - 2 * MARGIN_IN) * EMU), int(0.3 * EMU),
                    name="source")
    paragraph(tf, text, 9, color, first=True, font="Arial")


def bullets(slide, x, y, w, items, *, name, size=14, color=INK, gap=8):
    h = int((0.05 + len(items) * 0.5) * EMU)
    _, tf = textbox(slide, x, y, w, h, name=name)
    for i, txt in enumerate(items):
        paragraph(tf, txt, size, color, first=(i == 0), bullet=True, level=0,
                  font="Arial", space_after_pt=gap)


def stat_cards(slide, y, cards, *, accent_cycle=(AMBER, BLUE, TEAL, INK)):
    """Row of equal-width, equal-top, even-gutter KPI cards.

    Each card is a SINGLE filled shape whose own text frame holds the number +
    label as two paragraphs — so there is no separate textbox overlapping the
    card (check_pptx flags any two shapes whose boxes intersect).
    """
    n = len(cards)
    total_w = int((SW - 2 * MARGIN_IN) * EMU)
    gut = int(0.25 * EMU)
    cw = (total_w - (n - 1) * gut) // n
    chigh = int(1.7 * EMU)
    for i, (value, label) in enumerate(cards):
        x = M + i * (cw + gut)
        sp = rectangle(slide, x, y, cw, chigh, name=f"card_{i}", fill=CARD,
                       line=CARD_LINE, line_w_pt=0.75)
        tf = sp.text_frame
        tf.word_wrap = True
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        tf.margin_left = tf.margin_right = Pt(14)
        tf.margin_top = tf.margin_bottom = Pt(10)
        paragraph(tf, value, 30, accent_cycle[i % len(accent_cycle)], bold=True,
                  first=True, font="Georgia", space_after_pt=2)
        paragraph(tf, label, 11, MUTED, font="Arial", space_after_pt=0)
    return y + chigh


# --------------------------------------------------------------------------- #
# Slides
# --------------------------------------------------------------------------- #

# 1. COVER (dark)
s = new_slide(DARK)
_, tf = textbox(s, M, int(2.0 * EMU), int(11 * EMU), int(2.0 * EMU), name="title")
paragraph(tf, "The Financial Story of Take-Two Interactive", 44, LIGHTTX, bold=True,
          first=True, font="Georgia")
connector(s, M, int(3.95 * EMU), int(1.8 * EMU), AMBER, name="title_rule", weight_pt=3)
_, stf = textbox(s, M, int(4.15 * EMU), int(11 * EMU), int(0.9 * EMU), name="subtitle")
paragraph(stf, "Record fiscal 2021 results and the first half of fiscal 2022 — a "
          "publisher of owned, iconic franchises (Nasdaq: TTWO)", 16, FAINT, first=True, font="Arial")
paragraph(stf, "Fiscal year ended March 31, 2021, and the six months ended September 30, 2021",
          12, FAINT, italic=True, font="Arial", space_before_pt=8)
_, ftf = textbox(s, M, SRC_Y, int(11 * EMU), int(0.3 * EMU), name="footer")
paragraph(ftf, "Take-Two Interactive Software, Inc. · Nasdaq: TTWO", 10, AMBER, bold=True,
          first=True, font="Arial")

# 2. CONTEXT — two bulleted columns
s = new_slide(CREAM)
motif(s)
ct = title(s, "Context for this overview", accent=AMBER, ink=INK)
col_w = int(5.3 * EMU)
gut = int(0.4 * EMU)
right_x = M + col_w + gut
_, h1 = textbox(s, M, ct, col_w, int(0.4 * EMU), name="head_covers")
paragraph(h1, "What this covers", 16, INK, bold=True, first=True, font="Georgia")
_, h2 = textbox(s, right_x, ct, col_w, int(0.4 * EMU), name="head_built")
paragraph(h2, "How Take-Two is built", 16, INK, bold=True, first=True, font="Georgia")
by = ct + int(0.5 * EMU)
bullets(s, M, by, col_w, [
    "This overview summarizes Take-Two's financial performance across revenue, profitability, revenue mix, and capital position.",
    "It draws on three lenses: the multi-year growth trend, the FY2021 record year in detail, and the early-FY2022 trajectory.",
    "All figures are as reported in SEC filings; net bookings is the company's key operational (non-GAAP) metric, shown alongside GAAP net revenue.",
], name="body_covers", size=13)
bullets(s, right_x, by, col_w, [
    "Five publishing labels: Rockstar Games, 2K, Private Division, Social Point, and Playdots.",
    "27 proprietary franchises, internally owned across action, sports, strategy, and mobile genres.",
    "Grand Theft Auto is the cornerstone IP — GTA V alone has sold over 155 million units since launch.",
    "6,495 full-time employees as of March 31, 2021, with 77% in product development.",
], name="body_built", size=13)
source(s, "Source: Take-Two FY2021 Form 10-K; Q2 FY2022 earnings release & Form 10-Q.")

# 3. KPIS — four hero cards + supporting bullets
s = new_slide(CREAM)
motif(s)
ct = title(s, "Fiscal 2021 set company records across every headline metric", accent=AMBER, ink=INK)
cards_bottom = stat_cards(s, ct, [
    ("$3,373M", "Net Revenue"), ("$589M", "Net Income"),
    ("$5.09", "Diluted EPS"), ("$3,553M", "Net Bookings"),
])
bullets(s, M, cards_bottom + int(0.4 * EMU), int((SW - 2 * MARGIN_IN) * EMU), [
    "Net revenue rose 9% year over year to a company record $3.37 billion.",
    "Net income grew 46% to $589 million, the highest in company history.",
    "Net bookings — the forward-looking demand metric — reached $3.55 billion, up 19%.",
], name="body_kpis", size=14)
source(s, "Source: FY2021 10-K, Consolidated Statements of Operations; FY2021 earnings release (Net Bookings).")


def chart_slide(title_text, subtitle, categories, series, kind, palette, bullets_items,
                src, *, legend=True, chart_name="chart"):
    s = new_slide(CREAM)
    motif(s)
    ct = title(s, title_text, accent=AMBER, ink=INK, subtitle=subtitle)
    chart_w = int(7.0 * EMU)
    chart_h = int(4.3 * EMU)
    chart(s, M, ct, chart_w, chart_h, categories, series, palette=palette,
          name=chart_name, kind=kind, legend_bottom=legend)
    bx = M + chart_w + int(0.4 * EMU)
    bw = RIGHT - bx
    _, hh = textbox(s, bx, ct, bw, int(0.4 * EMU), name="head_sowhat")
    paragraph(hh, "Why it matters", 15, AMBER, bold=True, first=True, font="Georgia")
    bullets(s, bx, ct + int(0.5 * EMU), bw, bullets_items, name="body_sowhat", size=13)
    source(s, src)
    return s

# 4. REVENUE TREND — line chart
chart_slide(
    "Net revenue has grown every year, up 26% over three years to a record $3.4B",
    "Net revenue, fiscal years ended March 31 ($M)",
    ["FY2017", "FY2018", "FY2019", "FY2020", "FY2021"],
    [("Net revenue", [1780, 1793, 2668, 3089, 3373])],
    "line", [AMBER],
    ["Five straight years of top-line expansion, accelerating sharply from FY2019 onward.",
     "FY2021's $3.37B is an all-time high, +9% over FY2020 and +26% over FY2019.",
     "Growth has been driven by the shift to recurrent consumer spending (next section)."],
    "Source: FY2021 10-K (FY2019–FY2021 exact); stockanalysis.com (FY2017–FY2018). Net revenue is GAAP.",
    legend=False, chart_name="chart_revenue_trend")

# 5. PROFITABILITY — column chart, two series
chart_slide(
    "Margins expanded sharply: gross margin reached 54.5% and operating margin 18.7%",
    "Gross margin and operating margin, FY2019–FY2021 (% of net revenue)",
    ["FY2019", "FY2020", "FY2021"],
    [("Gross margin", [42.9, 50.1, 54.5]), ("Operating margin", [7.7, 13.8, 18.7])],
    "column", [AMBER, BLUE],
    ["Gross margin climbed nearly 12 points in two years as higher-margin digital and recurrent revenue grew.",
     "Income from operations more than tripled, from $207M (FY2019) to $629M (FY2021).",
     "Net income margin reached 17.5% in FY2021, up from 12.5% in FY2019."],
    "Source: FY2021 10-K, Statements of Operations. Margins computed on GAAP net revenue.",
    chart_name="chart_margins")

# 6. REVENUE MIX — column chart, two series
chart_slide(
    "Recurrent consumer spending became the majority of revenue, flipping the model",
    "Net revenue by type, FY2020 vs FY2021 ($M)",
    ["Recurrent consumer spending", "Full game and other"],
    [("FY2020", [1385, 1704]), ("FY2021", [2075, 1298])],
    "column", [BLUE, AMBER],
    ["RCS grew 50% year over year and became the majority of revenue for the first time.",
     "RCS reached 62% of net revenue in FY2021 (and 63% of net bookings), up from 45% a year earlier.",
     "The shift de-risks the model: revenue is increasingly recurring rather than launch-tied."],
    "Source: FY2021 10-K, net revenue by type. RCS = recurrent consumer spending.",
    chart_name="chart_mix")


def two_table_slide(title_text, left_cap, left_headers, left_rows, right_cap,
                    right_headers, right_rows, bullets_items, src, *, name_l, name_r,
                    lw=(1.6, 1.3, 1.0), rw=None):
    s = new_slide(CREAM)
    motif(s)
    ct = title(s, title_text, accent=AMBER, ink=INK)
    tbl_w = int(5.6 * EMU)
    gut = int(0.4 * EMU)
    right_x = M + tbl_w + gut
    cap_h = int(0.3 * EMU)
    # left caption + table
    _, lc = textbox(s, M, ct, tbl_w, cap_h, name=f"caption_{name_l}")
    paragraph(lc, left_cap, 11, MUTED, bold=True, first=True, font="Arial")
    ty = ct + int(0.34 * EMU)
    rows_h = (len(left_rows) + 1) * 0.4
    table(s, M, ty, tbl_w, int(rows_h * EMU), left_headers, left_rows,
          header_fill=INK, header_text=RGBColor(0xFF, 0xFF, 0xFF), ink=INK, name=name_l,
          col_widths_in=lw)
    # right caption + table
    _, rc = textbox(s, right_x, ct, tbl_w, cap_h, name=f"caption_{name_r}")
    paragraph(rc, right_cap, 11, MUTED, bold=True, first=True, font="Arial")
    table(s, right_x, ty, tbl_w, int(((len(right_rows) + 1) * 0.4) * EMU), right_headers,
          right_rows, header_fill=INK, header_text=RGBColor(0xFF, 0xFF, 0xFF), ink=INK,
          name=name_r, col_widths_in=rw)
    by = ty + int((rows_h + 0.5) * EMU)
    bullets(s, M, by, int((SW - 2 * MARGIN_IN) * EMU), bullets_items, name="body_takeaway", size=13)
    source(s, src)
    return s

# 7. CHANNEL & GEO — two tables
two_table_slide(
    "Revenue is overwhelmingly digital and led by the U.S., with 40% from abroad",
    "FY2021 net revenue by channel",
    ["Channel", "FY2021 ($M)", "Share"],
    [["Digital online", "2,919", "86.6%"], ["Physical retail & other", "453", "13.4%"]],
    "Net revenue by geography, FY2020–FY2021",
    ["Geography", "FY2020 ($M)", "FY2021 ($M)", "Growth"],
    [["United States", "1,776", "2,016", "+14%"], ["International", "1,313", "1,357", "+3%"]],
    ["Digital delivery reached 86.6% of net revenue (87% of net bookings), up from 63% two years earlier.",
     "The U.S. drove growth at +14%; international held 40% of revenue, with Asia ~7% of net bookings — a stated expansion opportunity."],
    "Source: FY2021 10-K, net revenue by distribution channel and by geography.",
    name_l="table_channel", name_r="table_geography",
    lw=(2.4, 1.7, 1.5), rw=(1.7, 1.4, 1.4, 1.1))

# 8. PLATFORM — pie chart
chart_slide(
    "Console anchors the business while PC and mobile broaden the base",
    "FY2021 net revenue by platform ($M)",
    ["Console", "PC and other"],
    [("FY2021", [2517, 856])],
    "pie", [AMBER, BLUE],
    ["Console delivered 75% of FY2021 net revenue, led by the GTA and NBA 2K franchises.",
     "PC and other contributed $856M (25%), including a growing mobile and PC-digital base.",
     "Grand Theft Auto products alone represented 29.2% of total net revenue."],
    "Source: FY2021 10-K, net revenue by platform.",
    legend=True, chart_name="chart_platform")

# 9. BALANCE SHEET — four cards
s = new_slide(CREAM)
motif(s)
ct = title(s, "A fortress balance sheet: $2.7B in cash and investments funds the pipeline",
           accent=AMBER, ink=INK)
cards_bottom = stat_cards(s, ct, [
    ("$1,423M", "Cash & Equivalents"), ("$1,309M", "Short-Term Investments"),
    ("$912M", "Operating Cash Flow"), ("$3,332M", "Stockholders' Equity"),
])
bullets(s, M, cards_bottom + int(0.4 * EMU), int((SW - 2 * MARGIN_IN) * EMU), [
    "Take-Two held $2.7 billion in cash and short-term investments at fiscal year-end, against modest liabilities.",
    "Operating cash flow rose 33% to $912 million, funding the multi-year slate internally.",
    "Total assets were $6.0 billion with $3.3 billion of stockholders' equity — ample capacity for content and M&A.",
], name="body_balance", size=14)
source(s, "Source: FY2021 10-K, Balance Sheet and Statement of Cash Flows.")

# 10. SECTION DIVIDER (dark)
s = new_slide(DARK)
_, tf = textbox(s, M, int(2.6 * EMU), int(11 * EMU), int(1.6 * EMU), name="section_label")
paragraph(tf, "Part Two: The first half of fiscal 2022", 40, LIGHTTX, bold=True, first=True, font="Georgia")
connector(s, M, int(4.0 * EMU), int(1.8 * EMU), AMBER, name="section_rule", weight_pt=3)
_, stf = textbox(s, M, int(4.2 * EMU), int(11 * EMU), int(0.6 * EMU), name="section_sub")
paragraph(stf, "Investing through a transition year while recurrent spending deepens", 16, FAINT, first=True, font="Arial")

# 11. Q2 FY22 — four cards
s = new_slide(CREAM)
motif(s)
ct = title(s, "Q2 FY2022 revenue held steady as the company invested for its largest-ever pipeline",
           accent=AMBER, ink=INK)
cards_bottom = stat_cards(s, ct, [
    ("$858M", "Q2 Net Revenue"), ("$20M", "Q2 Operating Income"),
    ("$0.09", "Q2 Diluted EPS"), ("$985M", "Q2 Net Bookings"),
])
bullets(s, M, cards_bottom + int(0.4 * EMU), int((SW - 2 * MARGIN_IN) * EMU), [
    "Net bookings grew 3% to $985M; recurrent consumer spending rose to 69% of net bookings.",
    "Operating income was $20M (down from $115M a year earlier) as Take-Two stepped up development investment.",
    "Gross margin was 46.8%; the quarter reflects a deliberate build-up ahead of a record forward slate.",
], name="body_q2", size=14)
source(s, "Source: Q2 FY2022 earnings release and Form 10-Q (quarter ended September 30, 2021).")

# 12. Q2 MIX — two tables (long title — must size for wrap)
two_table_slide(
    "In Q2, digital and recurrent spending deepened across both geography and channel",
    "Q2 FY22 net revenue by geography",
    ["Geography", "Q2 FY22 ($M)", "Share"],
    [["United States", "515", "60%"], ["International", "343", "40%"]],
    "Q2 FY22 net revenue by channel",
    ["Channel", "Q2 FY22 ($M)", "Share"],
    [["Digital online", "779", "91%"], ["Physical & other", "79", "9%"]],
    ["Digital reached 91% of quarterly net revenue (89% of net bookings) — the structural shift continues.",
     "The U.S./international split held near 60/40, consistent with the full-year FY2021 mix."],
    "Source: Q2 FY2022 earnings release (quarter ended September 30, 2021).",
    name_l="table_q2_geo", name_r="table_q2_channel",
    lw=(2.2, 2.0, 1.4), rw=(2.2, 2.0, 1.4))

# 13. MOBILE — column chart
chart_slide(
    "Mobile net revenue nearly doubled, led by Social Point's live-service titles",
    "Mobile net revenue, Q2 FY2021 vs Q2 FY2022 ($M)",
    ["Q2 FY2021", "Q2 FY2022"],
    [("Mobile net revenue", [61, 115])],
    "column", [TEAL],
    ["Mobile net revenue grew 89% year over year, the fastest-growing platform in the quarter.",
     "Growth was led by Social Point's live-service mobile portfolio and continued in-game spending.",
     "Mobile deepens the recurrent-revenue mix and broadens reach beyond console and PC."],
    "Source: Q2 FY2022 Form 10-Q. Mobile is reported within 'PC and other'.",
    legend=False, chart_name="chart_mobile")

# 14. QUOTE (dark) — motif top-right, no floating bars
s = new_slide(DARK)
cw, ch, gap = int(0.22 * EMU), int(0.1 * EMU), int(0.1 * EMU)
mx = RIGHT - 3 * cw - 2 * gap
for i, c in enumerate(CHIPS):
    rectangle(s, mx + i * (cw + gap), int(0.55 * EMU), cw, ch, name=f"motif_chip_{i}", fill=c)
_, tf = textbox(s, M, int(2.3 * EMU), int(11 * EMU), int(1.8 * EMU), name="quote")
paragraph(tf, "“We are entering our strongest development pipeline ever.”", 40,
          LIGHTTX, bold=True, first=True, font="Georgia")
connector(s, M, int(4.15 * EMU), int(1.6 * EMU), AMBER, name="quote_rule", weight_pt=3)
_, atf = textbox(s, M, int(4.35 * EMU), int(11 * EMU), int(0.6 * EMU), name="attribution")
paragraph(atf, "— Strauss Zelnick, Chairman & CEO, on Take-Two's multi-year slate", 16,
          FAINT, first=True, font="Arial")

# 15. CLOSING (dark) — takeaways. Use the title() helper (same as content
# slides) so the connector sits below the real wrap, not through it.
s = new_slide(DARK)
_ctop = title(s, "A record year and a deliberate investment phase",
              accent=AMBER, ink=LIGHTTX, name="title")
bullets(s, M, _ctop + int(0.15 * EMU), int(11 * EMU), [
    "FY2021 delivered record revenue ($3.37B), net income ($589M), and margins, on a model now majority-recurrent.",
    "The balance sheet — $2.7B in cash and investments, $912M operating cash flow — funds growth internally.",
    "H1 FY2022 reflects deliberate investment ahead of the strongest pipeline in company history.",
], name="body_closing", size=16, color=LIGHTTX, gap=12)
_, ftf = textbox(s, M, SRC_Y, int(11 * EMU), int(0.3 * EMU), name="footer")
paragraph(ftf, "Take-Two Interactive Software, Inc. · Nasdaq: TTWO", 11, AMBER, bold=True,
          first=True, font="Arial")

prs.save("ttwo_overview.pptx")
print("saved ttwo_overview.pptx ·", len(prs.slides._sldIdLst), "slides")
