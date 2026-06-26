from __future__ import annotations

from dataclasses import dataclass, field
from itertools import combinations

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE, PP_PLACEHOLDER

EMU_PER_INCH = 914400
TOL = 0.1
CAPTION_GAP = 0.5
TITLE_TYPES = {PP_PLACEHOLDER.TITLE, PP_PLACEHOLDER.CENTER_TITLE}
PRIMARY_KINDS = {"title", "body", "chart", "table", "picture", "text"}


def emu_to_in(value: int) -> float:
    """Convert a python-pptx length or raw EMU value to inches."""
    inches = value.inches if hasattr(value, "inches") else value / EMU_PER_INCH
    return round(inches, 3)


# --------------------------------------------------------------------------- #
# Shape collection
# --------------------------------------------------------------------------- #

@dataclass
class Box:
    """One leaf shape, with all geometry already in inches (slide coords)."""

    name: str
    kind: str          # title | body | chart | table | picture | text | other
    left: float
    top: float
    width: float
    height: float
    text: str = ""

    @property
    def right(self) -> float:
        return round(self.left + self.width, 3)

    @property
    def bottom(self) -> float:
        return round(self.top + self.height, 3)

    def as_dict(self) -> dict:
        return {
            "left": self.left,
            "top": self.top,
            "width": self.width,
            "height": self.height,
        }


@dataclass
class SlideInfo:
    index: int                       # 1-based, for human-facing messages
    boxes: list[Box] = field(default_factory=list)
    notes: list[str] = field(default_factory=list)


@dataclass
class DeckContext:
    slide_w: float                   # inches
    slide_h: float                   # inches
    slides: list[SlideInfo]


def _kind_of(shape, title_shape) -> str:
    """Classify a shape into one of our coarse kinds."""
    if shape is title_shape:
        return "title"
    # has_table / has_chart are safe on any shape and never raise.
    if getattr(shape, "has_table", False):
        return "table"
    if getattr(shape, "has_chart", False):
        return "chart"
    if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
        return "picture"
    if shape.is_placeholder:
        ph_type = shape.placeholder_format.type
        if ph_type in TITLE_TYPES:
            return "title"
        if ph_type == PP_PLACEHOLDER.BODY:
            return "body"
    if getattr(shape, "has_text_frame", False) and shape.has_text_frame:
        return "text"
    return "other"


def _group_transform(group):
    """Return the affine map (child coords -> parent coords) for a group.

    Reads off/ext (group box in parent coords) and chOff/chExt (child coord
    space) from the group's xfrm. Returns a closure mapping a child bbox
    (in EMU) to its bbox in the parent's coordinate space (in EMU).
    """
    xfrm = group._element.grpSpPr.xfrm
    ox, oy = xfrm.off.x, xfrm.off.y
    ecx, ecy = xfrm.ext.cx, xfrm.ext.cy
    cox, coy = xfrm.chOff.x, xfrm.chOff.y
    chcx, chcy = xfrm.chExt.cx, xfrm.chExt.cy

    # Guard degenerate child extents (would divide by zero).
    sx = (ecx / chcx) if chcx else 1.0
    sy = (ecy / chcy) if chcy else 1.0

    def to_parent(left, top, width, height):
        return (
            ox + (left - cox) * sx,
            oy + (top - coy) * sy,
            width * sx,
            height * sy,
        )

    return to_parent


def _walk(shape, title_shape, transforms, out: list[Box], notes: list[str]):
    """Recursively collect leaf shapes, applying group transforms.

    ``transforms`` is the chain of group maps from outermost to innermost; each
    leaf's EMU bbox is pushed back through every map (innermost first) to reach
    absolute slide coordinates.
    """
    if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
        to_parent = _group_transform(shape)
        for child in shape.shapes:
            _walk(child, None, transforms + [to_parent], out, notes)
        return

    left, top = shape.left, shape.top
    width, height = shape.width, shape.height
    if None in (left, top, width, height):
        # Inherited-geometry placeholder (no explicit position). Skip but note.
        notes.append(f"{shape.name or 'shape'}: inherited geometry (skipped)")
        return

    left, top, width, height = int(left), int(top), int(width), int(height)
    # Apply group maps from innermost outward to reach slide coordinates.
    for to_parent in reversed(transforms):
        left, top, width, height = to_parent(left, top, width, height)

    out.append(
        Box(
            name=shape.name or "shape",
            kind=_kind_of(shape, title_shape),
            left=emu_to_in(left),
            top=emu_to_in(top),
            width=emu_to_in(width),
            height=emu_to_in(height),
            text=(
                shape.text_frame.text.strip()
                if getattr(shape, "has_text_frame", False) and shape.has_text_frame
                else ""
            ),
        )
    )


def collect_shapes(slide, index: int) -> SlideInfo:
    """Flatten a slide into leaf Boxes in absolute slide inches."""
    info = SlideInfo(index=index + 1)
    title_shape = slide.shapes.title  # may be None
    for shape in slide.shapes:
        _walk(shape, title_shape, [], info.boxes, info.notes)
    return info


# --------------------------------------------------------------------------- #
# Problem helpers
# --------------------------------------------------------------------------- #

def _problem(criterion, slide, message, shapes=None, boxes=None) -> dict:
    return {
        "criterion": criterion,
        "slide": slide,
        "message": message,
        "shapes": shapes or [],
        "boxes_in": boxes or [],
    }


def _intersection_area(a: Box, b: Box) -> float:
    dx = min(a.right, b.right) - max(a.left, b.left)
    dy = min(a.bottom, b.bottom) - max(a.top, b.top)
    if dx > 0 and dy > 0:
        return dx * dy
    return 0.0


# --------------------------------------------------------------------------- #
# Checks (one per rubric criterion)
# --------------------------------------------------------------------------- #

def check_outer_margin_frame(ctx: DeckContext) -> list[dict]:
    crit = "layout-and-alignment-outer-margin-frame"
    out = []
    for s in ctx.slides:
        for b in s.boxes:
            off_slide = (
                b.left < 0 or b.top < 0
                or b.right > ctx.slide_w or b.bottom > ctx.slide_h
            )
            in_band = (
                b.left < TOL or b.top < TOL
                or b.right > ctx.slide_w - TOL or b.bottom > ctx.slide_h - TOL
            )
            if off_slide:
                out.append(_problem(
                    crit, s.index,
                    f"{b.name!r} bleeds off the slide",
                    [b.name], [b.as_dict()],
                ))
            elif in_band:
                out.append(_problem(
                    crit, s.index,
                    f"{b.name!r} intrudes into the {TOL} in outer margin band",
                    [b.name], [b.as_dict()],
                ))
    return out


def check_no_overlap(ctx: DeckContext) -> list[dict]:
    crit = "layout-and-alignment-no-overlap-or-collision"
    out = []
    for s in ctx.slides:
        for a, b in combinations(s.boxes, 2):
            area = _intersection_area(a, b)
            if area > 0:
                out.append(_problem(
                    crit, s.index,
                    f"{a.name!r} and {b.name!r} overlap by "
                    f"{round(area, 3)} in^2",
                    [a.name, b.name], [a.as_dict(), b.as_dict()],
                ))
    return out


def check_left_margin_consistency(ctx: DeckContext) -> list[dict]:
    crit = "layout-and-alignment-left-margin-consistency"
    out = []

    # Per-slide: dominant left edge = the most common block left edge.
    # We count *block* left edges, not every shape's left: peer columns within a
    # row share one block (the row), so only the row's leftmost edge counts.
    # Otherwise a legitimate 3-column layout would read as 3 distinct margins.
    slide_lefts = {}
    for s in ctx.slides:
        lefts = _block_left_edges(s.boxes)
        if not lefts:
            continue
        # Cluster lefts within TOL; dominant cluster's representative value.
        clusters = _cluster(lefts)
        # Distinct clusters per slide should be <= 2.
        if len(clusters) > 2:
            out.append(_problem(
                crit, s.index,
                f"slide uses {len(clusters)} distinct block left edges "
                f"(expected at most 2)",
            ))
        dominant = max(clusters, key=lambda c: len(c))
        slide_lefts[s.index] = sum(dominant) / len(dominant)

    if len(slide_lefts) >= 2:
        # Deck dominant left = most common slide-dominant value (within TOL).
        values = list(slide_lefts.values())
        deck_clusters = _cluster(values)
        deck_dominant_cluster = max(deck_clusters, key=lambda c: len(c))
        deck_left = sum(deck_dominant_cluster) / len(deck_dominant_cluster)
        deviating = [
            idx for idx, v in slide_lefts.items()
            if abs(v - deck_left) > TOL
        ]
        if len(deviating) >= 2:
            out.append(_problem(
                crit, deviating[0],
                f"{len(deviating)} slides (e.g. {sorted(deviating)}) use a left "
                f"edge differing by more than {TOL} in from the deck's dominant "
                f"left edge of {round(deck_left, 3)} in",
            ))
    return out


def check_edge_and_grid_alignment(ctx: DeckContext) -> list[dict]:
    crit = "layout-and-alignment-edge-and-grid-alignment"
    out = []
    for s in ctx.slides:
        for row in _rows(s.boxes):
            if len(row) < 2:
                continue
            widths = [b.width for b in row]
            tops = [b.top for b in row]
            if max(widths) - min(widths) > TOL:
                out.append(_problem(
                    crit, s.index,
                    f"peer boxes in a row have unequal widths "
                    f"({round(min(widths), 3)}..{round(max(widths), 3)} in)",
                    [b.name for b in row], [b.as_dict() for b in row],
                ))
            if max(tops) - min(tops) > TOL:
                out.append(_problem(
                    crit, s.index,
                    f"peer boxes in a row are not top-aligned "
                    f"({round(min(tops), 3)}..{round(max(tops), 3)} in)",
                    [b.name for b in row], [b.as_dict() for b in row],
                ))
    return out


def check_consistent_gutters(ctx: DeckContext) -> list[dict]:
    crit = "layout-and-alignment-consistent-gutters"
    out = []
    for s in ctx.slides:
        # Horizontal gutters within rows of 3+.
        for row in _rows(s.boxes):
            if len(row) < 3:
                continue
            ordered = sorted(row, key=lambda b: b.left)
            gaps = [
                ordered[i + 1].left - ordered[i].right
                for i in range(len(ordered) - 1)
            ]
            if gaps and max(gaps) - min(gaps) > TOL:
                out.append(_problem(
                    crit, s.index,
                    f"horizontal gutters in a row of {len(row)} vary "
                    f"({round(min(gaps), 3)}..{round(max(gaps), 3)} in)",
                    [b.name for b in ordered],
                    [b.as_dict() for b in ordered],
                ))
        # Vertical gutters within columns of 3+ genuine stacked peers.
        for col in _peer_columns(s.boxes):
            if len(col) < 3:
                continue
            ordered = sorted(col, key=lambda b: b.top)
            gaps = [
                ordered[i + 1].top - ordered[i].bottom
                for i in range(len(ordered) - 1)
            ]
            if gaps and max(gaps) - min(gaps) > TOL:
                out.append(_problem(
                    crit, s.index,
                    f"vertical gutters in a column of {len(col)} vary "
                    f"({round(min(gaps), 3)}..{round(max(gaps), 3)} in)",
                    [b.name for b in ordered],
                    [b.as_dict() for b in ordered],
                ))
    return out


def check_element_boxing(ctx: DeckContext) -> list[dict]:
    crit = "layout-and-alignment-element-boxing"
    out = []
    for s in ctx.slides:
        captions = [b for b in s.boxes if b.kind in ("text", "title") and b.text]
        for el in s.boxes:
            if el.kind not in ("chart", "table"):
                continue
            has_caption = any(
                # caption sits above the element, within CAPTION_GAP, x-overlapping
                (el.top - cap.bottom) >= -TOL
                and (el.top - cap.bottom) <= CAPTION_GAP
                and min(cap.right, el.right) - max(cap.left, el.left) > 0
                for cap in captions
            )
            if not has_caption:
                out.append(_problem(
                    crit, s.index,
                    f"{el.kind} {el.name!r} has no caption within "
                    f"{CAPTION_GAP} in above it",
                    [el.name], [el.as_dict()],
                ))
    return out


def check_vertical_rhythm(ctx: DeckContext) -> list[dict]:
    crit = "layout-and-alignment-vertical-rhythm"
    out = []
    title_tops = {}
    footer_tops = {}
    for s in ctx.slides:
        titles = [b for b in s.boxes if b.kind == "title"]
        if titles:
            title_tops[s.index] = min(b.top for b in titles)
        footers = [
            b for b in s.boxes
            if b.kind in ("text", "body") and b.top > ctx.slide_h - 1.0
        ]
        if footers:
            footer_tops[s.index] = max(footers, key=lambda b: b.top).top

    out += _rhythm_violation(crit, "title", title_tops)
    out += _rhythm_violation(crit, "footer/page-number", footer_tops)
    return out


# --------------------------------------------------------------------------- #
# Small geometry / clustering utilities
# --------------------------------------------------------------------------- #

def _cluster(values: list[float]) -> list[list[float]]:
    """Greedy 1-D clustering: values within TOL of a cluster's mean join it."""
    clusters: list[list[float]] = []
    for v in sorted(values):
        for c in clusters:
            if abs(v - (sum(c) / len(c))) <= TOL:
                c.append(v)
                break
        else:
            clusters.append([v])
    return clusters


# Titles span the slide and are not grid peers; exclude them from row/column
# peer detection so a title stacked above content isn't mistaken for a column.
PEER_KINDS = {"body", "chart", "table", "picture", "text"}


def _rows(boxes: list[Box]) -> list[list[Box]]:
    """Group boxes into peer rows.

    A peer row is 2+ boxes that share a top (within TOL) AND have similar
    heights (within TOL) -- i.e. a genuine multi-column band, not a title that
    happens to share a top with something.
    """
    peers = [b for b in boxes if b.kind in PEER_KINDS]
    rows: list[list[Box]] = []
    for b in sorted(peers, key=lambda x: x.top):
        for row in rows:
            if abs(b.top - row[0].top) <= TOL and abs(b.height - row[0].height) <= TOL:
                row.append(b)
                break
        else:
            rows.append([b])
    return [r for r in rows if len(r) >= 2]


def _peer_columns(boxes: list[Box]) -> list[list[Box]]:
    """Group boxes into peer columns.

    A peer column is 2+ boxes that share a left (within TOL) AND have similar
    widths (within TOL) -- a genuine stacked column, not coincidentally
    left-aligned heterogeneous shapes (e.g. a title above a chart).
    """
    peers = [b for b in boxes if b.kind in PEER_KINDS]
    cols: list[list[Box]] = []
    for b in sorted(peers, key=lambda x: x.left):
        for col in cols:
            if abs(b.left - col[0].left) <= TOL and abs(b.width - col[0].width) <= TOL:
                col.append(b)
                break
        else:
            cols.append([b])
    return [c for c in cols if len(c) >= 2]


def _block_left_edges(boxes: list[Box]) -> list[float]:
    """Left edges of block-level elements, one per peer row.

    Peer columns within a row collapse to the row's leftmost edge, so a clean
    N-column layout contributes a single block left edge rather than N.
    """
    rows = _rows(boxes)
    row_members = {id(b) for row in rows for b in row}
    lefts = []
    # Each peer row contributes only its leftmost edge.
    for row in rows:
        lefts.append(min(b.left for b in row))
    # Standalone primary elements (not part of any peer row) contribute their own.
    for b in boxes:
        if b.kind in PRIMARY_KINDS and id(b) not in row_members:
            lefts.append(b.left)
    return lefts


def _rhythm_violation(crit, label, tops: dict[int, float]) -> list[dict]:
    if len(tops) < 2:
        return []
    values = list(tops.values())
    if max(values) - min(values) > TOL:
        slides = sorted(tops.keys())
        return [_problem(
            crit, slides[0],
            f"{label} top y-offset varies across slides "
            f"({round(min(values), 3)}..{round(max(values), 3)} in)",
        )]
    return []


# --------------------------------------------------------------------------- #
# Entry point
# --------------------------------------------------------------------------- #

def check_pptx(path: str) -> dict:
    """Validate the layout of a saved .pptx and return structured problems.

    Returns a dict: {ok, problems, notes}. Each problem names the
    criterion, the 1-based slide number, the offending shapes, and their boxes
    (in inches), so the caller knows exactly what to move.

    NOTE: cannot detect text overflowing inside a box (python-pptx does not
    reflow text); only geometric problems are caught.
    """
    prs = Presentation(path)
    slides = [collect_shapes(s, i) for i, s in enumerate(prs.slides)]
    ctx = DeckContext(
        slide_w=emu_to_in(prs.slide_width),
        slide_h=emu_to_in(prs.slide_height),
        slides=slides,
    )

    problems: list[dict] = []
    problems += check_outer_margin_frame(ctx)
    problems += check_no_overlap(ctx)
    problems += check_left_margin_consistency(ctx)
    problems += check_edge_and_grid_alignment(ctx)
    problems += check_consistent_gutters(ctx)
    problems += check_element_boxing(ctx)
    problems += check_vertical_rhythm(ctx)

    notes = [
        {"slide": s.index, "notes": s.notes}
        for s in slides if s.notes
    ]
    return {
        "ok": not problems,
        "problems": problems,
        "notes": notes,
    }
