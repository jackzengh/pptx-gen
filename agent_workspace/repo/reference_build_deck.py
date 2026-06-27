#!/usr/bin/env python3
"""Build Take-Two Interactive FY2021 & H1 FY2022 financial overview deck."""

import os, sys
sys.path.insert(0, os.path.dirname(os.path.abspath(__file__)))

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR, MSO_AUTO_SIZE
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.chart.data import CategoryChartData
from pptx.oxml.ns import qn

# ==================== Constants ====================
EMU = 914400
SLIDE_W = 13.333
SLIDE_H = 7.5
MARGIN = 0.45
TITLE_TOP = 0.45
TITLE_PT = 28
SRC_Y = 6.9
PAGENUM_LEFT = 12.55
PAGENUM_W = 0.4
SRC_W = 11.5
CONTENT_W = SLIDE_W - 2 * MARGIN  # 12.433

# ==================== Palette ====================
NAVY = RGBColor(0x0B, 0x1F, 0x3A)
BLUE = RGBColor(0x1B, 0x4D, 0x7E)
RED = RGBColor(0xE6, 0x39, 0x46)
INK = RGBColor(0x1A, 0x23, 0x32)
MUTED = RGBColor(0x60, 0x6A, 0x72)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_TEXT = RGBColor(0xCA, 0xDC, 0xFC)
MUTED_LIGHT = RGBColor(0xA0, 0xAE, 0xC0)
GREY = RGBColor(0xB0, 0xB0, 0xB0)

# ==================== Skill Helpers (copied from python-pptx.md) ====================
_EMU_PER_POINT = 914400 / 72
_BULLET_GLYPHS = ["\u2022", "\u2013", "\u25e6"]
_HANG_EM = 1.1
_STEP_EM = 1.4

def bullet_indent_emu(level, size_pt):
    em = size_pt * _EMU_PER_POINT
    marL = round((_HANG_EM + _STEP_EM * max(0, level)) * em)
    indent = round(-_HANG_EM * em)
    return marL, indent

def _apply_bullet(p, level, size_pt, font):
    marL, indent = bullet_indent_emu(level, size_pt)
    pPr = p._p.get_or_add_pPr()
    pPr.set("marL", str(marL))
    pPr.set("indent", str(indent))
    for tag in ("a:buNone", "a:buChar", "a:buAutoNum", "a:buFont"):
        for el in pPr.findall(qn(tag)):
            pPr.remove(el)
    glyph = _BULLET_GLYPHS[min(max(0, level), len(_BULLET_GLYPHS) - 1)]
    pPr.append(pPr.makeelement(qn("a:buFont"), {"typeface": font}))
    pPr.append(pPr.makeelement(qn("a:buChar"), {"char": glyph}))

def style_run(run, size_pt, color, *, bold=False, italic=False, font="Arial"):
    run.font.size = Pt(size_pt)
    run.font.color.rgb = color
    run.font.bold = bold
    run.font.italic = italic
    run.font.name = font

def textbox(slide, left, top, width, height, *, name):
    tb = slide.shapes.add_textbox(Emu(left), Emu(top), Emu(width), Emu(height))
    tb.name = name
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Pt(0)
    tf.margin_top = tf.margin_bottom = Pt(0)
    return tb, tf

def paragraph(tf, text, size_pt, color, *, bold=False, italic=False, level=0,
              align=PP_ALIGN.LEFT, space_after_pt=6, space_before_pt=0,
              first=False, font="Arial", bullet=False):
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    p.level = level
    p.alignment = align
    p.space_after = Pt(space_after_pt)
    p.space_before = Pt(space_before_pt)
    if bullet:
        _apply_bullet(p, level, size_pt, font)
    run = p.add_run()
    run.text = text
    style_run(run, size_pt, color, bold=bold, italic=italic, font=font)
    return p

def rectangle(slide, left, top, width, height, *, name,
              fill=None, line=None, line_w_pt=0.75):
    sp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Emu(left), Emu(top), Emu(width), Emu(height))
    sp.name = name
    sp.shadow.inherit = False
    if fill is None:
        sp.fill.background()
    else:
        sp.fill.solid()
        sp.fill.fore_color.rgb = fill
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line
        sp.line.width = Pt(line_w_pt)
    return sp

def connector(slide, left, top, width, color, *, name, weight_pt=1.5):
    line = slide.shapes.add_connector(2, Emu(left), Emu(top), Emu(left + width), Emu(top))
    line.name = name
    line.line.color.rgb = color
    line.line.width = Pt(weight_pt)
    return line

# --- Higher-level helpers ---
MARGIN_IN = 0.45
TITLE_TOP_IN = 0.45
TITLE_GAP_IN = 0.14
CONN_TO_BODY_IN = 0.18

def title(slide, text, *, accent, ink, width_in=11.0, font="Georgia",
          name="title", italic=False, bold=True):
    from harness.text_metrics import measure_text_emu
    left = int(MARGIN_IN * EMU)
    top = int(TITLE_TOP_IN * EMU)
    w = int(width_in * EMU)
    _, h_emu = measure_text_emu(text, w, TITLE_PT, 1.15, font_family=font, bold=bold)
    h_emu = max(h_emu, int(0.7 * EMU))
    tb, tf = textbox(slide, left, top, w, h_emu, name=name)
    paragraph(tf, text, TITLE_PT, ink, bold=bold, italic=italic, first=True,
              font=font, space_after_pt=0)
    conn_y = top + h_emu + int(TITLE_GAP_IN * EMU)
    connector(slide, left, conn_y, int(1.6 * EMU), accent, name=f"{name}_rule", weight_pt=2.5)
    content_top = conn_y + int(CONN_TO_BODY_IN * EMU)
    return content_top

def chart(slide, left, top, width, height, categories, series, *, palette,
          name, kind="column", legend_bottom=True):
    data = CategoryChartData()
    data.categories = categories
    for label, vals in series:
        data.add_series(label, vals)
    xl = {"column": XL_CHART_TYPE.COLUMN_CLUSTERED, "bar": XL_CHART_TYPE.BAR_CLUSTERED,
          "line": XL_CHART_TYPE.LINE_MARKERS, "pie": XL_CHART_TYPE.PIE}[kind]
    gf = slide.shapes.add_chart(xl, Emu(left), Emu(top), Emu(width), Emu(height), data)
    gf.name = name
    ch = gf.chart
    ch.has_title = False
    if kind == "pie":
        ch.has_legend = True
        ch.legend.position = XL_LEGEND_POSITION.BOTTOM
        ch.legend.include_in_layout = False
        for i, pt in enumerate(ch.plots[0].series[0].points):
            pt.format.fill.solid()
            pt.format.fill.fore_color.rgb = palette[i % len(palette)]
        return gf
    ch.has_legend = bool(legend_bottom) and len(series) > 1
    if ch.has_legend:
        ch.legend.position = XL_LEGEND_POSITION.BOTTOM
        ch.legend.include_in_layout = False
    for i, plot_series in enumerate(ch.plots[0].series):
        plot_series.format.fill.solid()
        plot_series.format.fill.fore_color.rgb = palette[i % len(palette)]
    return gf

def table(slide, left, top, width, height, headers, rows, *, header_fill, header_text,
          ink, name, col_widths_in=None, header_pt=13, body_pt=13, row_h_in=0.4):
    gf = slide.shapes.add_table(len(rows) + 1, len(headers), Emu(left), Emu(top),
                                Emu(width), Emu(height))
    gf.name = name
    t = gf.table
    for j, htext in enumerate(headers):
        c = t.cell(0, j)
        c.fill.solid(); c.fill.fore_color.rgb = header_fill
        c.text = str(htext)
        p = c.text_frame.paragraphs[0]
        p.runs[0].font.size = Pt(header_pt); p.runs[0].font.bold = True
        p.runs[0].font.color.rgb = header_text; p.runs[0].font.name = "Arial"
    for i, row in enumerate(rows, start=1):
        for j, val in enumerate(row):
            c = t.cell(i, j)
            c.text = str(val)
            r = c.text_frame.paragraphs[0].runs[0]
            r.font.size = Pt(body_pt); r.font.color.rgb = ink; r.font.name = "Arial"
    t.rows[0].height = Emu(int(row_h_in * EMU))
    for i in range(1, len(rows) + 1):
        t.rows[i].height = Emu(int(row_h_in * EMU))
    if col_widths_in:
        for j, w_in in enumerate(col_widths_in):
            t.columns[j].width = Emu(int(w_in * EMU))
    return gf

# ==================== Custom Helpers ====================

def format_chart(gf, muted_color, num_format=None):
    ch = gf.chart
    try:
        ch.category_axis.tick_labels.font.size = Pt(10)
        ch.category_axis.tick_labels.font.color.rgb = muted_color
        ch.category_axis.tick_labels.font.name = "Arial"
        ch.value_axis.tick_labels.font.size = Pt(10)
        ch.value_axis.tick_labels.font.color.rgb = muted_color
        ch.value_axis.tick_labels.font.name = "Arial"
        if num_format:
            ch.value_axis.tick_labels.number_format = num_format
    except Exception:
        pass
    if ch.has_legend:
        ch.legend.font.size = Pt(10)
        ch.legend.font.color.rgb = muted_color
        ch.legend.font.name = "Arial"
    return gf

def stat_card(slide, left, top, width, height, name, value, label, fill_color, text_color):
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                  Emu(left), Emu(top), Emu(width), Emu(height))
    card.name = name
    card.shadow.inherit = False
    card.fill.solid()
    card.fill.fore_color.rgb = fill_color
    card.line.fill.background()
    tf = card.text_frame
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.NONE
    tf.margin_left = Pt(6); tf.margin_right = Pt(6)
    tf.margin_top = Pt(8); tf.margin_bottom = Pt(8)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p1 = tf.paragraphs[0]
    p1.alignment = PP_ALIGN.CENTER
    r1 = p1.add_run()
    r1.text = value
    style_run(r1, 30, text_color, bold=True, font="Arial")
    p1.space_after = Pt(4)
    p2 = tf.add_paragraph()
    p2.alignment = PP_ALIGN.CENTER
    r2 = p2.add_run()
    r2.text = label
    style_run(r2, 11, text_color, font="Arial")
    return card

def source_line(slide, text, color, y=SRC_Y):
    left = int(MARGIN * EMU)
    top = int(y * EMU)
    w = int(SRC_W * EMU)
    h = int(0.2 * EMU)
    tb, tf = textbox(slide, left, top, w, h, name="meta_source")
    paragraph(tf, text, 9, color, first=True, space_after_pt=0)
    return tb

def page_number(slide, num, color, y=SRC_Y):
    left = int(PAGENUM_LEFT * EMU)
    top = int(y * EMU)
    w = int(PAGENUM_W * EMU)
    h = int(0.2 * EMU)
    tb, tf = textbox(slide, left, top, w, h, name=f"meta_pagenum_{num}")
    paragraph(tf, str(num), 9, color, first=True, align=PP_ALIGN.RIGHT, space_after_pt=0)
    return tb

def dark_bg(slide):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = NAVY

def light_bg(slide):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = WHITE

def bullet_mixed(tf, segments, size_pt, color, level=0, first=False, space_after_pt=6):
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    p.level = level
    p.alignment = PP_ALIGN.LEFT
    p.space_after = Pt(space_after_pt)
    _apply_bullet(p, level, size_pt, "Arial")
    for text, bold in segments:
        run = p.add_run()
        run.text = text
        style_run(run, size_pt, color, bold=bold, font="Arial")
    return p

def caption_box(slide, left, top, width, text, color, name):
    h = int(0.25 * EMU)
    tb, tf = textbox(slide, left, top, width, h, name=name)
    paragraph(tf, text, 9, color, first=True, space_after_pt=0)
    return tb

# ==================== Card geometry ====================
CARD_W = int(2.92 * EMU)
CARD_H = int(1.6 * EMU)
CARD_GUTTER = int(0.25 * EMU)
CARD_X1 = int(0.45 * EMU)
CARD_X2 = CARD_X1 + CARD_W + CARD_GUTTER
CARD_X3 = CARD_X2 + CARD_W + CARD_GUTTER
CARD_X4 = CARD_X3 + CARD_W + CARD_GUTTER

# ==================== Slide Builders ====================

def add_slide(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])

def build_cards(slide, content_top, cards):
    for i, (val, label, name) in enumerate(cards):
        x = [CARD_X1, CARD_X2, CARD_X3, CARD_X4][i]
        stat_card(slide, x, content_top, CARD_W, CARD_H, name, val, label, BLUE, WHITE)

def build_bullets(slide, top, left, width, height, bullets, color, name, space_after=8):
    tb, tf = textbox(slide, left, top, width, height, name=name)
    for i, item in enumerate(bullets):
        if isinstance(item, list):
            bullet_mixed(tf, item, 13, color, first=(i==0), space_after_pt=space_after)
        else:
            paragraph(tf, item, 13, color, bullet=True, first=(i==0), space_after_pt=space_after)
    return tb


# ---- Slide 1: Cover ----
def slide_cover(prs):
    slide = add_slide(prs)
    dark_bg(slide)
    ct = title(slide, "The Financial Story of Take-Two Interactive",
               accent=RED, ink=WHITE)
    # body text
    _, tf = textbox(slide, int(0.45*EMU), ct, int(11.0*EMU), int(0.6*EMU), name="body_cover")
    paragraph(tf, "Record fiscal 2021 results and the first half of fiscal 2022 \u2014 a publisher of owned, iconic franchises (Nasdaq: TTWO)",
              13, LIGHT_TEXT, first=True, space_after_pt=0)
    # subtitle near bottom
    _, tf2 = textbox(slide, int(0.45*EMU), int(6.3*EMU), int(10.0*EMU), int(0.3*EMU), name="body_subtitle")
    paragraph(tf2, "Fiscal year ended March 31, 2021, and the six months ended September 30, 2021",
              11, MUTED_LIGHT, first=True, space_after_pt=0)
    # accent bar at bottom
    rectangle(slide, 0, int(7.35*EMU), int(13.333*EMU), int(0.15*EMU),
              name="background_accent_bar", fill=RED)


# ---- Slide 2: Context ----
def slide_context(prs):
    slide = add_slide(prs)
    light_bg(slide)
    ct = title(slide, "Context for this overview", accent=RED, ink=INK)
    # left column
    _, tfl = textbox(slide, int(0.45*EMU), ct, int(6.0*EMU), int(4.8*EMU), name="body_left_col")
    paragraph(tfl, "What this covers", 13, BLUE, bold=True, first=True, space_after_pt=8)
    paragraph(tfl, "This overview summarizes Take-Two's financial performance across revenue, profitability, revenue mix, and capital position.", 13, INK, bullet=True, space_after_pt=6)
    paragraph(tfl, "It draws on three lenses: the multi-year growth trend, the FY2021 record year in detail, and the early-FY2022 trajectory.", 13, INK, bullet=True, space_after_pt=6)
    paragraph(tfl, "All figures are as reported in SEC filings; net bookings is the company's key operational (non-GAAP) metric and is shown alongside GAAP net revenue.", 13, INK, bullet=True, space_after_pt=0)
    # right column
    _, tfr = textbox(slide, int(6.75*EMU), ct, int(6.0*EMU), int(4.8*EMU), name="body_right_col")
    paragraph(tfr, "How Take-Two is built", 13, BLUE, bold=True, first=True, space_after_pt=8)
    bullet_mixed(tfr, [("Five publishing labels:", True), (" Rockstar Games, 2K, Private Division, Social Point, and Playdots.", False)], 13, INK, space_after_pt=6)
    bullet_mixed(tfr, [("27 proprietary franchises", True), (", internally owned across action, sports, strategy, and mobile genres.", False)], 13, INK, space_after_pt=6)
    bullet_mixed(tfr, [("Grand Theft Auto", True), (" is the cornerstone IP \u2014 GTA V alone has sold in over 155 million units since launch.", False)], 13, INK, space_after_pt=6)
    bullet_mixed(tfr, [("6,495 full-time employees", True), (" as of March 31, 2021, with 77% in product development.", False)], 13, INK, space_after_pt=0)
    source_line(slide, "Source: Take-Two FY2021 Form 10-K; Q2 FY2022 earnings release & Form 10-Q.", MUTED)
    page_number(slide, 2, MUTED)


# ---- Slide 3: KPIs ----
def slide_kpis(prs):
    slide = add_slide(prs)
    light_bg(slide)
    ct = title(slide, "Fiscal 2021 set company records across every headline metric", accent=RED, ink=INK)
    cards = [
        ("$3,373M", "Net Revenue", "card_net_revenue"),
        ("$589M", "Net Income", "card_net_income"),
        ("$5.09", "Diluted EPS", "card_diluted_eps"),
        ("$3,553M", "Net Bookings", "card_net_bookings"),
    ]
    build_cards(slide, ct, cards)
    bt = ct + CARD_H + int(0.3*EMU)
    bullets = [
        "Net revenue rose 9% year over year to a company record $3.37 billion.",
        "Net income grew 46% to $589 million, the highest in company history.",
        "Net bookings \u2014 the forward-looking demand metric \u2014 reached $3.55 billion, up 19%.",
    ]
    build_bullets(slide, bt, int(0.45*EMU), int(CONTENT_W*EMU), int(2.5*EMU), bullets, INK, "body_kpi_bullets")
    source_line(slide, "Source: FY2021 10-K, Consolidated Statements of Operations; FY2021 earnings release (Net Bookings).", MUTED)
    page_number(slide, 3, MUTED)


# ---- Slide 4: Revenue Trend ----
def slide_revenue_trend(prs):
    slide = add_slide(prs)
    light_bg(slide)
    ct = title(slide, "Net revenue has grown every year, up 26% over three years to a record $3.4B", accent=RED, ink=INK)
    cap_top = ct
    chart_top = ct + int(0.25*EMU)
    chart_w = int(7.5*EMU)
    chart_h = int(4.5*EMU)
    caption_box(slide, int(0.45*EMU), cap_top, chart_w, "Net revenue, fiscal years ended March 31 ($M)", MUTED, "caption_chart_revenue")
    gf = chart(slide, int(0.45*EMU), chart_top, chart_w, chart_h,
               ["FY2017", "FY2018", "FY2019", "FY2020", "FY2021"],
               [("Net revenue", [1780, 1793, 2668, 3089, 3373])],
               palette=[BLUE], name="chart_revenue_trend", kind="line")
    format_chart(gf, MUTED, "#,##0")
    bullets = [
        "Five straight years of top-line expansion, accelerating sharply from FY2019 onward.",
        "FY2021's $3.37B is an all-time high, +9% over FY2020 and +26% over FY2019.",
        "Growth has been driven by the shift to recurrent consumer spending (see next section).",
    ]
    build_bullets(slide, ct, int(8.3*EMU), int(4.55*EMU), int(4.5*EMU), bullets, INK, "body_rev_trend_bullets")
    source_line(slide, "Source: FY2021 10-K (FY2019\u2013FY2021 exact); stockanalysis.com (FY2017\u2013FY2018). Note: net revenue is GAAP; differs from net bookings.", MUTED)
    page_number(slide, 4, MUTED)


# ---- Slide 5: Profitability ----
def slide_profitability(prs):
    slide = add_slide(prs)
    light_bg(slide)
    ct = title(slide, "Margins expanded sharply: gross margin reached 54.5% and operating margin 18.7%", accent=RED, ink=INK)
    chart_w = int(7.5*EMU)
    chart_h = int(4.5*EMU)
    chart_top = ct + int(0.25*EMU)
    caption_box(slide, int(0.45*EMU), ct, chart_w, "Gross margin and operating margin, FY2019\u2013FY2021 (% of net revenue)", MUTED, "caption_chart_margins")
    gf = chart(slide, int(0.45*EMU), chart_top, chart_w, chart_h,
               ["FY2019", "FY2020", "FY2021"],
               [("Gross margin", [42.9, 50.1, 54.5]), ("Operating margin", [7.7, 13.8, 18.7])],
               palette=[BLUE, RED], name="chart_profitability", kind="column")
    format_chart(gf, MUTED, "0.0")
    bullets = [
        "Gross margin climbed nearly 12 points in two years as higher-margin digital and recurrent revenue grew.",
        "Income from operations more than tripled, from $207M (FY2019) to $629M (FY2021).",
        "Net income margin reached 17.5% in FY2021, up from 12.5% in FY2019.",
    ]
    build_bullets(slide, ct, int(8.3*EMU), int(4.55*EMU), int(4.5*EMU), bullets, INK, "body_profit_bullets")
    source_line(slide, "Source: FY2021 10-K, Statements of Operations. Margins computed on GAAP net revenue.", MUTED)
    page_number(slide, 5, MUTED)


# ---- Slide 6: Revenue Mix ----
def slide_revenue_mix(prs):
    slide = add_slide(prs)
    light_bg(slide)
    ct = title(slide, "Recurrent consumer spending became the majority of revenue, flipping the model", accent=RED, ink=INK)
    chart_w = int(7.5*EMU)
    chart_h = int(4.5*EMU)
    chart_top = ct + int(0.25*EMU)
    caption_box(slide, int(0.45*EMU), ct, chart_w, "Net revenue by type, FY2020 vs FY2021 ($M)", MUTED, "caption_chart_mix")
    gf = chart(slide, int(0.45*EMU), chart_top, chart_w, chart_h,
               ["FY2020", "FY2021"],
               [("Recurrent consumer spending", [1385, 2075]), ("Full game and other", [1704, 1298])],
               palette=[RED, GREY], name="chart_revenue_mix", kind="column")
    format_chart(gf, MUTED, "#,##0")
    bullets = [
        "RCS grew 50% year over year and became the majority of revenue for the first time.",
        "RCS reached 62% of net revenue in FY2021 (and 63% of net bookings), up from 45% a year earlier.",
        "The shift de-risks the model: revenue is increasingly recurring rather than tied to new-title launches.",
    ]
    build_bullets(slide, ct, int(8.3*EMU), int(4.55*EMU), int(4.5*EMU), bullets, INK, "body_mix_bullets")
    source_line(slide, "Source: FY2021 10-K, net revenue by type. RCS = recurrent consumer spending (in-game purchases, virtual currency, add-on content).", MUTED)
    page_number(slide, 6, MUTED)


# ---- Slide 7: Channel / Geo ----
def slide_channel_geo(prs):
    slide = add_slide(prs)
    light_bg(slide)
    ct = title(slide, "Revenue is overwhelmingly digital and led by the U.S., with 40% from abroad", accent=RED, ink=INK)
    tbl_w = int(6.0*EMU)
    tbl_h = int(1.2*EMU)
    tbl_top = ct + int(0.25*EMU)
    # left table caption + table
    caption_box(slide, int(0.45*EMU), ct, tbl_w, "Net revenue by distribution channel", MUTED, "caption_table_channel")
    table(slide, int(0.45*EMU), tbl_top, tbl_w, tbl_h,
          ["Channel", "FY2021 ($M)", "Share"],
          [["Digital online", "2,919", "86.6%"], ["Physical retail & other", "453", "13.4%"]],
          header_fill=BLUE, header_text=WHITE, ink=INK, name="table_channel",
          col_widths_in=[2.8, 1.8, 1.4])
    # right table caption + table
    caption_box(slide, int(6.75*EMU), ct, tbl_w, "Net revenue by geography", MUTED, "caption_table_geo")
    table(slide, int(6.75*EMU), tbl_top, tbl_w, tbl_h,
          ["Geography", "FY2020 ($M)", "FY2021 ($M)", "Growth"],
          [["United States", "1,776", "2,016", "+14%"], ["International", "1,313", "1,357", "+3%"]],
          header_fill=BLUE, header_text=WHITE, ink=INK, name="table_geo",
          col_widths_in=[1.8, 1.4, 1.4, 1.4])
    bt = tbl_top + tbl_h + int(0.3*EMU)
    bullets = [
        "Digital delivery reached 86.6% of net revenue (87% of net bookings), up from 63% just two years earlier.",
        "The U.S. drove growth at +14%; international held 40% of revenue, with Asia ~7% of net bookings \u2014 a stated expansion opportunity.",
    ]
    build_bullets(slide, bt, int(0.45*EMU), int(CONTENT_W*EMU), int(2.5*EMU), bullets, INK, "body_channel_geo_bullets")
    source_line(slide, "Source: FY2021 10-K, net revenue by distribution channel and by geography.", MUTED)
    page_number(slide, 7, MUTED)


# ---- Slide 8: Platform ----
def slide_platform(prs):
    slide = add_slide(prs)
    light_bg(slide)
    ct = title(slide, "Console anchors the business while PC and mobile broaden the base", accent=RED, ink=INK)
    chart_w = int(7.5*EMU)
    chart_h = int(4.5*EMU)
    chart_top = ct + int(0.25*EMU)
    caption_box(slide, int(0.45*EMU), ct, chart_w, "FY2021 net revenue by platform ($M)", MUTED, "caption_chart_platform")
    gf = chart(slide, int(0.45*EMU), chart_top, chart_w, chart_h,
               ["Console", "PC and other"],
               [("FY2021", [2517, 856])],
               palette=[BLUE, RED], name="chart_platform", kind="pie")
    format_chart(gf, MUTED)
    bullets = [
        "Console delivered 75% of FY2021 net revenue, led by the GTA and NBA 2K franchises.",
        "PC and other contributed $856M (25%), including a growing mobile and PC-digital base.",
        "Grand Theft Auto products alone represented 29.2% of total net revenue.",
    ]
    build_bullets(slide, ct, int(8.3*EMU), int(4.55*EMU), int(4.5*EMU), bullets, INK, "body_platform_bullets")
    source_line(slide, "Source: FY2021 10-K, net revenue by platform.", MUTED)
    page_number(slide, 8, MUTED)


# ---- Slide 9: Balance Sheet ----
def slide_balance_sheet(prs):
    slide = add_slide(prs)
    light_bg(slide)
    ct = title(slide, "A fortress balance sheet: $2.7B in cash and investments funds the pipeline", accent=RED, ink=INK)
    cards = [
        ("$1,423M", "Cash & Equivalents", "card_cash"),
        ("$1,309M", "Short-Term Investments", "card_sti"),
        ("$912M", "Operating Cash Flow", "card_ocf"),
        ("$3,332M", "Stockholders' Equity", "card_equity"),
    ]
    build_cards(slide, ct, cards)
    bt = ct + CARD_H + int(0.3*EMU)
    bullets = [
        "Take-Two held $2.7 billion in cash and short-term investments at fiscal year-end, against modest liabilities.",
        "Operating cash flow rose 33% to $912 million, funding development of the multi-year slate internally.",
        "Total assets were $6.0 billion with $3.3 billion of stockholders' equity \u2014 ample capacity for content investment and M&A.",
    ]
    build_bullets(slide, bt, int(0.45*EMU), int(CONTENT_W*EMU), int(2.5*EMU), bullets, INK, "body_balance_bullets")
    source_line(slide, "Source: FY2021 10-K, Balance Sheet and Statement of Cash Flows.", MUTED)
    page_number(slide, 9, MUTED)


# ---- Slide 10: Section Divider ----
def slide_section(prs):
    slide = add_slide(prs)
    dark_bg(slide)
    ct = title(slide, "Part Two: The first half of fiscal 2022", accent=RED, ink=WHITE)
    _, tf = textbox(slide, int(0.45*EMU), ct, int(11.0*EMU), int(0.5*EMU), name="body_section")
    paragraph(tf, "Investing through a transition year while recurrent spending deepens",
              13, LIGHT_TEXT, first=True, space_after_pt=0)
    rectangle(slide, 0, int(7.35*EMU), int(13.333*EMU), int(0.15*EMU),
              name="background_accent_bar", fill=RED)
    page_number(slide, 10, MUTED_LIGHT)


# ---- Slide 11: Q2 FY2022 ----
def slide_q2fy22(prs):
    slide = add_slide(prs)
    light_bg(slide)
    ct = title(slide, "Q2 FY2022 revenue held steady as the company invested for its largest-ever pipeline", accent=RED, ink=INK)
    cards = [
        ("$858M", "Q2 Net Revenue", "card_q2_rev"),
        ("$20M", "Q2 Operating Income", "card_q2_oi"),
        ("$0.09", "Q2 Diluted EPS", "card_q2_eps"),
        ("$985M", "Q2 Net Bookings", "card_q2_bookings"),
    ]
    build_cards(slide, ct, cards)
    bt = ct + CARD_H + int(0.3*EMU)
    bullets = [
        "Net bookings grew 3% to $985M; recurrent consumer spending rose to 69% of net bookings.",
        "Operating income was $20M (down from $115M a year earlier) as Take-Two stepped up investment in development.",
        "Gross margin was 46.8%; the quarter reflects a deliberate build-up ahead of a record forward slate.",
    ]
    build_bullets(slide, bt, int(0.45*EMU), int(CONTENT_W*EMU), int(2.5*EMU), bullets, INK, "body_q2_bullets")
    source_line(slide, "Source: Q2 FY2022 earnings release and Form 10-Q (quarter ended September 30, 2021).", MUTED)
    page_number(slide, 11, MUTED)


# ---- Slide 12: Q2 Mix ----
def slide_q2_mix(prs):
    slide = add_slide(prs)
    light_bg(slide)
    ct = title(slide, "In Q2, digital and recurrent spending deepened across both geography and channel", accent=RED, ink=INK)
    tbl_w = int(6.0*EMU)
    tbl_h = int(1.2*EMU)
    tbl_top = ct + int(0.25*EMU)
    caption_box(slide, int(0.45*EMU), ct, tbl_w, "Q2 net revenue by geography", MUTED, "caption_table_q2_geo")
    table(slide, int(0.45*EMU), tbl_top, tbl_w, tbl_h,
          ["Geography", "Q2 FY22 ($M)", "Share"],
          [["United States", "515", "60%"], ["International", "343", "40%"]],
          header_fill=BLUE, header_text=WHITE, ink=INK, name="table_q2_geo",
          col_widths_in=[2.8, 1.8, 1.4])
    caption_box(slide, int(6.75*EMU), ct, tbl_w, "Q2 net revenue by channel", MUTED, "caption_table_q2_channel")
    table(slide, int(6.75*EMU), tbl_top, tbl_w, tbl_h,
          ["Channel", "Q2 FY22 ($M)", "Share"],
          [["Digital online", "779", "91%"], ["Physical & other", "79", "9%"]],
          header_fill=BLUE, header_text=WHITE, ink=INK, name="table_q2_channel",
          col_widths_in=[2.8, 1.8, 1.4])
    bt = tbl_top + tbl_h + int(0.3*EMU)
    bullets = [
        "Digital reached 91% of quarterly net revenue (89% of net bookings) \u2014 the structural shift continues.",
        "The U.S./international split held near 60/40, consistent with the full-year FY2021 mix.",
    ]
    build_bullets(slide, bt, int(0.45*EMU), int(CONTENT_W*EMU), int(2.5*EMU), bullets, INK, "body_q2_mix_bullets")
    source_line(slide, "Source: Q2 FY2022 earnings release (quarter ended September 30, 2021).", MUTED)
    page_number(slide, 12, MUTED)


# ---- Slide 13: Mobile ----
def slide_mobile(prs):
    slide = add_slide(prs)
    light_bg(slide)
    ct = title(slide, "Mobile net revenue nearly doubled, led by Social Point's live-service titles", accent=RED, ink=INK)
    chart_w = int(7.5*EMU)
    chart_h = int(4.5*EMU)
    chart_top = ct + int(0.25*EMU)
    caption_box(slide, int(0.45*EMU), ct, chart_w, "Mobile net revenue, Q2 FY2021 vs Q2 FY2022 ($M)", MUTED, "caption_chart_mobile")
    gf = chart(slide, int(0.45*EMU), chart_top, chart_w, chart_h,
               ["Q2 FY2021", "Q2 FY2022"],
               [("Mobile net revenue", [61, 115])],
               palette=[RED], name="chart_mobile", kind="column")
    format_chart(gf, MUTED, "#,##0")
    bullets = [
        "Mobile net revenue grew 89% year over year, the fastest-growing platform in the quarter.",
        "Growth was led by Social Point's live-service mobile portfolio and continued in-game spending.",
        "Mobile deepens the recurrent-revenue mix and broadens reach beyond console and PC.",
    ]
    build_bullets(slide, ct, int(8.3*EMU), int(4.55*EMU), int(4.5*EMU), bullets, INK, "body_mobile_bullets")
    source_line(slide, "Source: Q2 FY2022 Form 10-Q. Mobile is reported within \u201cPC and other.\u201d", MUTED)
    page_number(slide, 13, MUTED)


# ---- Slide 14: Quote ----
def slide_quote(prs):
    slide = add_slide(prs)
    dark_bg(slide)
    ct = title(slide, "\u201cWe are entering our strongest development pipeline ever.\u201d",
               accent=RED, ink=WHITE, italic=True, bold=False, width_in=11.5)
    _, tf = textbox(slide, int(0.45*EMU), ct, int(11.0*EMU), int(0.5*EMU), name="body_quote_attr")
    paragraph(tf, "\u2014 Strauss Zelnick, Chairman & CEO, on Take-Two's multi-year slate",
              13, LIGHT_TEXT, first=True, space_after_pt=0)
    rectangle(slide, 0, int(7.35*EMU), int(13.333*EMU), int(0.15*EMU),
              name="background_accent_bar", fill=RED)
    page_number(slide, 14, MUTED_LIGHT)


# ---- Slide 15: Closing ----
def slide_closing(prs):
    slide = add_slide(prs)
    dark_bg(slide)
    ct = title(slide, "A record year and a deliberate investment phase", accent=RED, ink=WHITE)
    bullets = [
        "FY2021 delivered record revenue ($3.37B), net income ($589M), and margins, on a model now majority-recurrent.",
        "The balance sheet \u2014 $2.7B in cash and investments, $912M operating cash flow \u2014 funds growth internally.",
        "H1 FY2022 reflects deliberate investment ahead of the strongest pipeline in company history.",
    ]
    build_bullets(slide, ct, int(0.45*EMU), int(11.5*EMU), int(3.0*EMU), bullets, LIGHT_TEXT, "body_closing_bullets")
    _, tf = textbox(slide, int(0.45*EMU), int(6.0*EMU), int(11.0*EMU), int(0.3*EMU), name="body_company")
    paragraph(tf, "Take-Two Interactive Software, Inc. \u00b7 Nasdaq: TTWO",
              11, WHITE, first=True, space_after_pt=0)
    source_line(slide, "Source: FY2021 10-K; Q2 FY2022 earnings release.", MUTED_LIGHT)
    page_number(slide, 15, MUTED_LIGHT)


# ==================== Main ====================
def main():
    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W)
    prs.slide_height = Inches(SLIDE_H)

    slide_cover(prs)       # 1
    slide_context(prs)     # 2
    slide_kpis(prs)        # 3
    slide_revenue_trend(prs)  # 4
    slide_profitability(prs)  # 5
    slide_revenue_mix(prs)    # 6
    slide_channel_geo(prs)    # 7
    slide_platform(prs)       # 8
    slide_balance_sheet(prs)  # 9
    slide_section(prs)        # 10
    slide_q2fy22(prs)         # 11
    slide_q2_mix(prs)         # 12
    slide_mobile(prs)         # 13
    slide_quote(prs)          # 14
    slide_closing(prs)        # 15

    out = os.path.join(os.path.dirname(os.path.abspath(__file__)), "ttwo_overview.pptx")
    prs.save(out)
    print(f"Saved {out} with {len(prs.slides)} slides")

if __name__ == "__main__":
    main()
